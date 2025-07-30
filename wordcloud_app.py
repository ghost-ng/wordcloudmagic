import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from ttkbootstrap.toast import ToastNotification
from ttkbootstrap.dialogs.colorchooser import ColorChooserDialog
from ttkbootstrap.widgets import Meter, Floodgauge
from tkinter import filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
import tkinter as tk
import tkinter.font as tkFont
import os
import sys
import threading
from PIL import Image, ImageTk, ImageDraw, ImageFont
import numpy as np
import platform
import subprocess
import json
from wordcloud import WordCloud, STOPWORDS
# Set matplotlib backend BEFORE importing pyplot
import matplotlib
matplotlib.use('TkAgg')

import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.colors import LinearSegmentedColormap
plt.ioff()  # Turn off interactive mode to prevent popup windows

# File handling imports
import PyPDF2
from docx import Document
from pptx import Presentation
import re
from datetime import datetime
from __version__ import __version__

def get_resource_path(relative_path):
    # For config and log files, use current working directory when running as PyInstaller app
    if hasattr(sys, '_MEIPASS') and (relative_path.startswith('configs') or relative_path.startswith('logs')):
        # Running as PyInstaller bundle - use current working directory
        base_path = os.getcwd()
        return os.path.join(base_path, relative_path)
    
    # For other resources (templates, assets, etc)
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.dirname(os.path.abspath(__file__))

    return os.path.join(base_path, relative_path)

# Debug print function
def debug_print(msg):
    """Print debug messages"""
    print(f"[DEBUG] {msg}")

class ToastManager:
    """Manages stacked toast notifications"""
    def __init__(self, root):
        self.root = root
        self.active_toasts = []
        self.toast_gap = 10
        self.base_y_offset = 50
        self.screen_width = root.winfo_screenwidth()
    
    def wrap_text(self, text, max_width=50):
        """Manually wrap text by inserting newlines"""
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 > max_width:
                lines.append(' '.join(current_line))
                current_line = [word]
                current_length = len(word)
            else:
                current_line.append(word)
                current_length += len(word) + 1
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)
    
    def show_toast(self, message, style="info", duration=15000):
        """Show a stacked toast notification"""
        # Wrap text for better display
        wrapped_message = self.wrap_text(message)
        
        # Calculate Y position based on existing toasts
        y_position = self.base_y_offset
        
        if self.active_toasts:
            # Calculate position based on actual heights of existing toasts
            for toast_data in self.active_toasts:
                if toast_data['toast'].toplevel and toast_data['toast'].toplevel.winfo_exists():
                    # Get actual height after toast is displayed
                    try:
                        toast_height = toast_data['toast'].toplevel.winfo_height()
                        # If height is 1, toast hasn't been rendered yet, use estimate
                        if toast_height == 1:
                            toast_height = toast_data.get('estimated_height', 80)
                        y_position += toast_height + self.toast_gap
                    except:
                        y_position += 80 + self.toast_gap  # Fallback height
        
        # Create toast with proper icon
        toast = ToastNotification(
            title="WordCloud Magic",
            message=wrapped_message,
            duration=duration,
            bootstyle=style,
            icon="✅" if style == "success" else "⚠" if style == "warning" else "✗" if style in ["danger", "error"] else "ℹ"
        )
        
        # Calculate X position (right side of screen with margin)
        x_position = self.screen_width - 350  # 350 is approximate toast width
        
        # Store toast data with estimated height
        toast_data = {
            'toast': toast,
            'y_position': y_position,
            'estimated_height': 80 + (message.count('\n') * 20)  # Estimate based on lines
        }
        self.active_toasts.append(toast_data)
        
        # Show the toast at calculated position
        toast.show_toast()
        
        # Override the position after showing
        def position_toast():
            if toast.toplevel and toast.toplevel.winfo_exists():
                toast.toplevel.geometry(f"+{x_position}+{y_position}")
                # Update with actual height once rendered
                toast_data['actual_height'] = toast.toplevel.winfo_height()
        
        # Position after a short delay to ensure toast is created
        self.root.after(10, position_toast)
        
        # Schedule removal from tracking after duration
        def remove_toast():
            # Remove this toast from active list
            self.active_toasts = [t for t in self.active_toasts if t['toast'] != toast]
            self._reposition_toasts()
        
        # Schedule removal slightly after the toast duration
        self.root.after(duration + 100, remove_toast)
    
    def _reposition_toasts(self):
        """Reposition remaining toasts after one is removed"""
        y_position = self.base_y_offset
        x_position = self.screen_width - 350
        
        for toast_data in self.active_toasts:
            toast = toast_data['toast']
            if hasattr(toast, 'toplevel') and toast.toplevel and toast.toplevel.winfo_exists():
                try:
                    toast.toplevel.geometry(f"+{x_position}+{y_position}")
                    toast_height = toast.toplevel.winfo_height()
                    if toast_height == 1:  # Not rendered yet
                        toast_height = toast_data.get('estimated_height', 80)
                    y_position += toast_height + self.toast_gap
                except:
                    pass

class FontListbox(ttk.Frame):
    """Custom font selector that displays fonts in their actual style"""
    def __init__(self, master, font_dict, textvariable=None, width=35, height=6, **kwargs):
        super().__init__(master, **kwargs)
        self.font_dict = font_dict
        self.textvariable = textvariable
        self.fonts_loaded = {}
        self.selected_index = -1
        self.items = []
        
        # Create frame for the selector
        self.columnconfigure(0, weight=1)
        self.rowconfigure(0, weight=1)
        
        # Create scrollbar
        scrollbar = ttk.Scrollbar(self, orient="vertical", bootstyle="primary-round")
        scrollbar.grid(row=0, column=1, sticky=(N, S))
        
        # Create Canvas
        self.canvas = tk.Canvas(self, 
                               width=width * 8,  # Approximate width in pixels
                               height=height * 22,  # Approximate height in pixels
                               bg='white',
                               highlightthickness=1)
        self.canvas.grid(row=0, column=0, sticky=(N, S, E, W))
        scrollbar.config(command=self.canvas.yview)
        self.canvas.config(yscrollcommand=scrollbar.set)
        
        # Bind events
        self.canvas.bind('<Button-1>', self._on_click)
        self.canvas.bind('<Configure>', self._on_canvas_configure)
        self.canvas.bind('<MouseWheel>', self._on_mousewheel)
        
        # Populate fonts
        self._populate_fonts()
        
    def _populate_fonts(self):
        """Populate the canvas with fonts in their actual styles"""
        # Clear existing items
        self.canvas.delete("all")
        self.items = []
        self.fonts_loaded = {}
        
        y_position = 5
        item_height = 25
        
        for i, font_name in enumerate(sorted(self.font_dict.keys())):
            # Try to create font
            try:
                font_face = self.font_dict[font_name]
                item_font = tkFont.Font(family=font_face, size=12)
                self.fonts_loaded[font_name] = item_font
            except:
                # If font fails to load, use default
                item_font = tkFont.Font(family='Segoe UI', size=12)
            
            # Create text item
            text_id = self.canvas.create_text(10, y_position + item_height//2,
                                             text=font_name,
                                             font=item_font,
                                             anchor='w',
                                             fill='black',
                                             tags=f"font_{i}")
            
            # Create selection rectangle (initially hidden) with outline instead of fill
            rect_id = self.canvas.create_rectangle(2, y_position, 
                                                  self.canvas.winfo_width() - 2, 
                                                  y_position + item_height,
                                                  fill='#e1f0ff',
                                                  outline='#0078d4',
                                                  width=2,
                                                  state='hidden',
                                                  tags=f"select_{i}")
            
            self.items.append({
                'name': font_name,
                'text_id': text_id,
                'rect_id': rect_id,
                'y_start': y_position,
                'y_end': y_position + item_height
            })
            
            y_position += item_height
        
        # Update scroll region
        self.canvas.config(scrollregion=self.canvas.bbox("all"))
        
        # Set initial selection
        if self.textvariable:
            current_value = self.textvariable.get()
            if current_value in self.font_dict:
                index = sorted(self.font_dict.keys()).index(current_value)
                self._select_item(index)
    
    def _on_click(self, event):
        """Handle click events"""
        # Convert canvas coordinates
        canvas_y = self.canvas.canvasy(event.y)
        
        # Find clicked item
        for i, item in enumerate(self.items):
            if item['y_start'] <= canvas_y <= item['y_end']:
                self._select_item(i)
                break
    
    def _select_item(self, index):
        """Select an item by index"""
        # Deselect previous
        if 0 <= self.selected_index < len(self.items):
            prev_item = self.items[self.selected_index]
            self.canvas.itemconfig(prev_item['rect_id'], state='hidden')
            self.canvas.itemconfig(prev_item['text_id'], fill='black')
            # Restore normal font
            if prev_item['name'] in self.fonts_loaded:
                font = self.fonts_loaded[prev_item['name']]
                font.configure(weight='normal')
        
        # Select new item
        if 0 <= index < len(self.items):
            self.selected_index = index
            item = self.items[index]
            # Show selection rectangle behind text
            self.canvas.tag_lower(item['rect_id'])  # Put rectangle behind text
            self.canvas.itemconfig(item['rect_id'], state='normal')
            self.canvas.itemconfig(item['text_id'], fill='#0078d4')  # Blue text
            
            # Make font bold for selected item
            if item['name'] in self.fonts_loaded:
                font = self.fonts_loaded[item['name']]
                font.configure(weight='bold')
            
            # Update variable
            if self.textvariable:
                self.textvariable.set(item['name'])
                self.event_generate('<<FontSelected>>', when='tail')
            
            # Ensure item is visible
            bbox = self.canvas.bbox(item['text_id'])
            if bbox:
                self.canvas.yview_moveto(bbox[1] / self.canvas.winfo_height())
    
    def _on_canvas_configure(self, event):
        """Update rectangles when canvas is resized"""
        canvas_width = event.width
        for item in self.items:
            self.canvas.coords(item['rect_id'], 
                             2, item['y_start'], 
                             canvas_width - 2, item['y_end'])
    
    def _on_mousewheel(self, event):
        """Handle mouse wheel scrolling"""
        self.canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
    
    def set_fonts(self, font_dict):
        """Update the available fonts"""
        self.font_dict = font_dict
        self._populate_fonts()

class ModernWordCloudApp:
    # Application version
    VERSION = __version__
    
    def print_debug(self, message):
        """Print debug message if in debug mode"""
        if hasattr(self, 'debug_mode') and self.debug_mode:
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            debug_msg = f"[DEBUG] {timestamp} - {message}"
            print(debug_msg)
            
            # Also write to log file if logging is enabled
            if hasattr(self, 'log_file') and self.log_file:
                try:
                    with open(self.log_file, 'a', encoding='utf-8') as f:
                        f.write(debug_msg + '\n')
                except Exception as e:
                    print(f"[ERROR] Failed to write to log: {e}")
    
    def print_info(self, message):
        """Print info message"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        info_msg = f"[INFO] {timestamp} - {message}"
        print(info_msg)
        
        # Also log if in debug mode
        if hasattr(self, 'debug_mode') and self.debug_mode and hasattr(self, 'log_file') and self.log_file:
            try:
                with open(self.log_file, 'a', encoding='utf-8') as f:
                    f.write(info_msg + '\n')
            except:
                pass
    
    def print_warning(self, message):
        """Print warning message"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        warn_msg = f"[WARNING] {timestamp} - {message}"
        print(warn_msg)
        
        # Also log if in debug mode
        if hasattr(self, 'debug_mode') and self.debug_mode and hasattr(self, 'log_file') and self.log_file:
            try:
                with open(self.log_file, 'a', encoding='utf-8') as f:
                    f.write(warn_msg + '\n')
            except:
                pass
    
    def print_fail(self, message):
        """Print failure/error message"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        error_msg = f"[ERROR] {timestamp} - {message}"
        print(error_msg)
        
        # Also log if in debug mode
        if hasattr(self, 'debug_mode') and self.debug_mode and hasattr(self, 'log_file') and self.log_file:
            try:
                with open(self.log_file, 'a', encoding='utf-8') as f:
                    f.write(error_msg + '\n')
            except:
                pass
    
    def init_logging(self):
        """Initialize logging to file"""
        try:
            # Create logs directory
            log_dir = get_resource_path('logs')
            os.makedirs(log_dir, exist_ok=True)
            
            # Create log file with timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            self.log_file = os.path.join(log_dir, f'wordcloud_debug_{timestamp}.log')
            
            # Write header
            with open(self.log_file, 'w', encoding='utf-8') as f:
                f.write(f"WordCloud Magic Debug Log\n")
                f.write(f"Started: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"{'='*60}\n\n")
            
            self.print_info(f"Logging to: {self.log_file}")
        except Exception as e:
            print(f"[ERROR] Failed to initialize logging: {e}")
            self.log_file = None
    
    def toggle_debug_mode(self):
        """Toggle debug mode on/off"""
        self.debug_mode = self.debug_var.get()
        
        if self.debug_mode:
            if not self.log_file:
                self.init_logging()
            self.show_toast("Debug mode enabled - logging to file", "info")
            self.print_debug("Debug mode toggled ON")
        else:
            self.show_toast("Debug mode disabled", "info")
            self.print_debug("Debug mode toggled OFF")
            # Note: we keep the log file open for the session
    
    def create_custom_gradients(self):
        """Create and register custom color gradients"""
        gradients = {}
        
        # Sunset Sky - Orange → Pink → Purple
        sunset_colors = ['#FF8C00', '#FF69B4', '#8B008B']
        gradients['sunset_sky'] = LinearSegmentedColormap.from_list('sunset_sky', sunset_colors)
        
        # Deep Ocean - Deep Blue → Teal → Light Blue
        ocean_colors = ['#000080', '#008B8B', '#87CEEB']
        gradients['deep_ocean'] = LinearSegmentedColormap.from_list('deep_ocean', ocean_colors)
        
        # Forest - Dark Green → Green → Light Green
        forest_colors = ['#006400', '#228B22', '#90EE90']
        gradients['forest'] = LinearSegmentedColormap.from_list('forest', forest_colors)
        
        # Fire - Red → Orange → Yellow
        fire_colors = ['#DC143C', '#FF8C00', '#FFD700']
        gradients['fire'] = LinearSegmentedColormap.from_list('fire', fire_colors)
        
        # Cotton Candy - Pink → Light Blue → Lavender
        cotton_colors = ['#FFB6C1', '#87CEFA', '#E6E6FA']
        gradients['cotton_candy'] = LinearSegmentedColormap.from_list('cotton_candy', cotton_colors)
        
        # Fall Leaves - Brown → Orange → Gold
        fall_colors = ['#8B4513', '#FF8C00', '#FFD700']
        gradients['fall_leaves'] = LinearSegmentedColormap.from_list('fall_leaves', fall_colors)
        
        # Berry - Deep Purple → Magenta → Pink
        berry_colors = ['#4B0082', '#FF00FF', '#FFC0CB']
        gradients['berry'] = LinearSegmentedColormap.from_list('berry', berry_colors)
        
        # Mint - Dark Teal → Mint → White
        mint_colors = ['#008080', '#98FB98', '#FFFFFF']
        gradients['mint'] = LinearSegmentedColormap.from_list('mint', mint_colors)
        
        # Volcano - Black → Red → Orange → Yellow
        volcano_colors = ['#000000', '#8B0000', '#FF4500', '#FFFF00']
        gradients['volcano'] = LinearSegmentedColormap.from_list('volcano', volcano_colors)
        
        # Aurora (Northern Lights) - Dark Blue → Green → Purple → Pink
        aurora_colors = ['#191970', '#00FF00', '#9370DB', '#FF1493']
        gradients['aurora'] = LinearSegmentedColormap.from_list('aurora', aurora_colors)
        
        # Hacker - Lime Green → Black
        hacker_colors = ['#00FF00', '#00AA00', '#005500', '#000000']
        gradients['hacker'] = LinearSegmentedColormap.from_list('hacker', hacker_colors)
        
        # Solarized Dark
        solarized_dark_colors = ['#002b36', '#073642', '#586e75', '#657b83', '#839496', '#93a1a1']
        gradients['solarized_dark'] = LinearSegmentedColormap.from_list('solarized_dark', solarized_dark_colors)
        
        # Solarized Light
        solarized_light_colors = ['#fdf6e3', '#eee8d5', '#93a1a1', '#839496', '#657b83', '#586e75']
        gradients['solarized_light'] = LinearSegmentedColormap.from_list('solarized_light', solarized_light_colors)
        
        # Rose Pine
        rose_pine_colors = ['#191724', '#1f1d2e', '#403d52', '#e0def4', '#eb6f92', '#f6c177']
        gradients['rose_pine'] = LinearSegmentedColormap.from_list('rose_pine', rose_pine_colors)
        
        # Grape - Deep Purple → Light Purple
        grape_colors = ['#2D1B69', '#512DA8', '#7E57C2', '#AB47BC', '#CE93D8']
        gradients['grape'] = LinearSegmentedColormap.from_list('grape', grape_colors)
        
        # Dracula
        dracula_colors = ['#282a36', '#44475a', '#6272a4', '#bd93f9', '#ff79c6', '#f8f8f2']
        gradients['dracula'] = LinearSegmentedColormap.from_list('dracula', dracula_colors)
        
        # Gruvbox
        gruvbox_colors = ['#282828', '#3c3836', '#504945', '#928374', '#d5c4a1', '#fbf1c7']
        gradients['gruvbox'] = LinearSegmentedColormap.from_list('gruvbox', gruvbox_colors)
        
        # Monokai
        monokai_colors = ['#272822', '#49483e', '#75715e', '#a6e22e', '#f92672', '#66d9ef']
        gradients['monokai'] = LinearSegmentedColormap.from_list('monokai', monokai_colors)
        
        # Army - Military Greens
        army_colors = ['#4B5320', '#556B2F', '#6B8E23', '#8FBC8F', '#90EE90']
        gradients['army'] = LinearSegmentedColormap.from_list('army', army_colors)
        
        # Air Force - Sky Blues
        airforce_colors = ['#00308F', '#0047AB', '#4169E1', '#6495ED', '#87CEEB']
        gradients['airforce'] = LinearSegmentedColormap.from_list('airforce', airforce_colors)
        
        # Cyber - Neon Cyan → Dark
        cyber_colors = ['#000000', '#0D0D0D', "#A6FF00", "#00D10A", '#1E90FF']
        gradients['cyber'] = LinearSegmentedColormap.from_list('cyber', cyber_colors)
        
        # Navy - Deep Ocean Blues
        navy_colors = ['#000080', '#002FA7', '#003F87', '#1560BD', '#4682B4']
        gradients['navy'] = LinearSegmentedColormap.from_list('navy', navy_colors)
        
        # Volcano - Fiery reds, oranges, yellows
        volcano_colors = ['#310600', '#950a11', '#f06625', '#f5b91d', '#f7f002']
        gradients['volcano'] = LinearSegmentedColormap.from_list('volcano', volcano_colors)
        
        # Lilac - Soft purples and pinks
        lilac_colors = ['#896790', '#B69CCF', '#D7ABE6', '#E7D1FF', '#F9EDFD']
        gradients['lilac'] = LinearSegmentedColormap.from_list('lilac', lilac_colors)
        
        # Cyberpunk - Neon pink, blue, purple
        cyberpunk_colors = ['#091833', '#133e7c', '#711c91', '#ea00d9', '#0abdc6']
        gradients['cyberpunk'] = LinearSegmentedColormap.from_list('cyberpunk', cyberpunk_colors)
        
        # Tron - Blue, cyan, orange
        tron_colors = ['#030504', '#062474', '#0EF8F8', '#7DFDFE', '#F4AF2D']
        gradients['tron'] = LinearSegmentedColormap.from_list('tron', tron_colors)
        
        # The Grid - Dark grey, neon green, blue
        grid_colors = ['#1A1A1A', '#333333', '#39FF14', '#03D8F3', '#00FFFF']
        gradients['grid'] = LinearSegmentedColormap.from_list('grid', grid_colors)
        
        # Fiber - Bright blue, magenta, cyan, purple
        fiber_colors = ['#0000FF', '#0080FF', '#00FFFF', '#8080FF', '#FF00FF']
        gradients['fiber'] = LinearSegmentedColormap.from_list('fiber', fiber_colors)
        
        # Register all custom colormaps with matplotlib
        for name, cmap in gradients.items():
            matplotlib.colormaps.register(cmap, name=name)
        
        return gradients
    
    def __init__(self, root):
        self.root = root
        self.root.title(f"WordCloud Magic v{self.VERSION} - Modern Word Cloud Generator")
        self.root.geometry("1300x850")
        self.root.state('zoomed')  # Start maximized

        # Asset paths - handle PyInstaller bundle
        if hasattr(sys, '_MEIPASS'):
            self.assets_dir = os.path.join(sys._MEIPASS, 'assets')
        else:
            self.assets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')
        self.icons = {}
        self.load_assets()
        
        # Initialize debug mode from command line arguments
        self.debug_mode = '--debug' in sys.argv
        self.log_file = None
        
        # Initialize logging if debug mode
        if self.debug_mode:
            self.init_logging()
            self.print_info("Debug mode enabled")
        
        # Log startup mode and version
        self.print_info(f"WordCloud Magic v{self.VERSION} starting...")
        if hasattr(sys, '_MEIPASS'):
            self.print_info("Running in EXE mode (PyInstaller bundle)")
            self.print_info(f"Bundle location: {sys._MEIPASS}")
            self.print_info(f"Working directory: {os.getcwd()}")
        else:
            self.print_info("Running in SCRIPT mode (Python interpreter)")
            self.print_info(f"Script location: {os.path.dirname(os.path.abspath(__file__))}")
        
        # Initialize toast manager
        self.toast_manager = ToastManager(self.root)
        
        # Flag to track UI readiness
        self.ui_ready = False
        
        # Available themes
        self.light_themes = [
            "cosmo", "flatly", "litera", "minty", "lumen", 
            "sandstone", "yeti", "pulse", "united", "morph",
            "journal", "simplex", "cerculean"
        ]
        self.dark_themes = [
            "darkly", "superhero", "solar", "cyborg", "vapor"
        ]
        self.themes = self.light_themes  # Start with light themes
        self.current_theme = tk.StringVar(value="cosmo")
        
        # Variables
        self.working_folder = tk.StringVar(value="No folder selected")
        self.text_content = ""
        self.mask_image = None  # For backward compatibility
        self.image_mask_image = None  # Store image mask separately
        self.text_mask_image = None   # Store text mask separately
        self.mask_path = tk.StringVar(value="No mask selected")
        self.min_word_length = tk.IntVar(value=3)
        self.max_word_length = tk.IntVar(value=20)
        self.forbidden_words = set()  # Start empty, will be populated from text area
        self.selected_colormap = "viridis"
        self.color_mode = tk.StringVar(value="preset")  # "single", "preset", or "custom"
        self.single_color = tk.StringVar(value="#0078D4")
        self.custom_gradient_colors = ["#FF0000", "#00FF00", "#0000FF"]  # Default RGB
        self.toast = ToastNotification(
            title="WordCloud Magic",
            message="",
            duration=3000,
            bootstyle=SUCCESS
        )
        
        # Text mask variables
        self.text_mask_input = tk.StringVar(value="")
        self.text_mask_font_size = tk.IntVar(value=200)
        self.text_mask_bold = tk.BooleanVar(value=True)
        self.text_mask_italic = tk.BooleanVar(value=False)
        self.text_mask_words_per_line = tk.IntVar(value=1)  # Words per line for multi-line text
        self.text_mask_font = tk.StringVar(value="Arial Black")  # Selected font
        
        # Available fonts for text mask
        self.available_fonts = {
            "Arial Black": "Arial Black",
            "Impact": "Impact",
            "Arial": "Arial", 
            "Helvetica": "Helvetica",
            "Times New Roman": "Times New Roman",
            "Georgia": "Georgia",
            "Verdana": "Verdana",
            "Comic Sans MS": "Comic Sans MS",
            "Trebuchet MS": "Trebuchet MS",
            "Courier New": "Courier New",
            "Calibri": "Calibri",
            "Cambria": "Cambria",
            "Tahoma": "Tahoma",
            "Century Gothic": "Century Gothic",
            "Palatino": "Palatino Linotype"
        }
        
        # Canvas settings
        self.canvas_width = tk.IntVar(value=800)
        self.canvas_height = tk.IntVar(value=600)
        self.bg_color = tk.StringVar(value="#FFFFFF")
        self.lock_aspect_ratio = tk.BooleanVar(value=False)
        self.aspect_ratio = 800 / 600  # Initial aspect ratio
        
        # Preview scale setting
        self.preview_scale = tk.IntVar(value=100)  # Default to 100%
        
        # Bind canvas size changes to preview update
        self.canvas_width.trace('w', self.update_preview_size)
        self.canvas_height.trace('w', self.update_preview_size)
        
        # Outline settings
        self.outline_width = tk.IntVar(value=0)
        self.outline_color = tk.StringVar(value="#000000")
        self.outline_widgets = []  # Keep track of outline widgets
        
        # Word orientation and mode
        self.prefer_horizontal = tk.DoubleVar(value=0.9)
        self.rgba_mode = tk.BooleanVar(value=False)
        self.show_transparency = tk.BooleanVar(value=True)
        self.max_words = tk.IntVar(value=200)
        self.scale = tk.IntVar(value=1)
        
        # Create custom gradients
        self.custom_gradients = self.create_custom_gradients()
        
        # Color schemes with descriptions
        self.color_schemes = {
            # Military/Service themes first
            "Army": "army",
            "Air Force": "airforce",
            "Navy": "navy",
            "Cyber": "cyber",
            # New colorful themes
            "Volcano": "volcano",
            "Lilac": "lilac", 
            "Cyberpunk": "cyberpunk",
            "Tron": "tron",
            "The Grid": "grid",
            "Fiber": "fiber",
            # Standard colormaps
            "Viridis": "viridis",
            "Plasma": "plasma",
            "Inferno": "inferno",
            "Magma": "magma",
            "Cool": "cool",
            "Hot": "hot",
            "Spring": "spring",
            "Summer": "summer",
            "Autumn": "autumn",
            "Winter": "winter",
            "Ocean": "ocean",
            "Rainbow": "rainbow",
            "Sunset": "RdYlBu",
            "Pastel": "Pastel1",
            "Dark": "Dark2",
            "Paired": "Paired",
            # New custom gradients
            "Sunset Sky": "sunset_sky",
            "Deep Ocean": "deep_ocean",
            "Forest": "forest",
            "Fire": "fire",
            "Cotton Candy": "cotton_candy",
            "Fall Leaves": "fall_leaves",
            "Berry": "berry",
            "Mint": "mint",
            "Volcano": "volcano",
            "Aurora": "aurora",
            "Neon": "neon",
            "Mystic": "mystic",
            "Hacker": "hacker",
            "SolarizedDk": "solarized_dark",
            "SolarizedLt": "solarized_light",
            "Rose Pine": "rose_pine",
            "Grape": "grape",
            "Dracula": "dracula",
            "Gruvbox": "gruvbox",
            "Monokai": "monokai"
        }
        
        # Initialize dark_mode before loading theme preference
        self.dark_mode = tk.BooleanVar(value=False)
        
        # Load theme preference before creating UI
        self.load_theme_preference()
        
        self.create_ui()
        
        # Mark UI as ready
        self.ui_ready = True
        
        # Initialize status bar labels
        self.update_color_scheme_label()
        
        # Bind window close event
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        # Auto-load configuration if exists (after UI is created)
        def load_config_and_theme():
            self.auto_load_config()
            # Always load theme from theme.json after config to ensure it takes precedence
            self.load_theme_preference()
        
        self.root.after(100, load_config_and_theme)
        
        # Validate available fonts after UI creation (in a thread to avoid blocking)
        threading.Thread(target=self.validate_fonts, daemon=True).start()
    
    def create_menu(self):
        """Create the menu bar"""
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)
        
        # File menu
        file_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="File", menu=file_menu)
        
        file_menu.add_command(label="Load Config", command=self.import_config)
        file_menu.add_command(label="Save Config As...", command=self.export_config)
        file_menu.add_command(label="Save Config", command=self.save_config_locally)
        file_menu.add_separator()
        file_menu.add_command(label="Reset", command=self.reset_app)
        file_menu.add_separator()
        file_menu.add_command(label="Help", command=self.show_help)
        file_menu.add_command(label=f"About (v{self.VERSION})", command=self.show_about)
        file_menu.add_separator()
        file_menu.add_command(label="Exit", command=self.on_closing)
        
    def create_ui(self):
        """Create the main UI"""
        # Create menu bar
        self.create_menu()
        
        # Top bar for theme selection and messages
        top_bar = ttk.Frame(self.root)
        top_bar.pack(fill=X, padx=10, pady=(10, 5))
        
        # Create message bar on the left side of top bar
        self.create_message_bar(top_bar)
        
        # Theme selector on the right
        theme_frame = ttk.Frame(top_bar)
        theme_frame.pack(side=RIGHT)
        
        # Dark mode toggle (already initialized in __init__)
        dark_mode_check = ttk.Checkbutton(theme_frame, 
                                         text="🌙 Dark Mode",
                                         variable=self.dark_mode,
                                         command=self.toggle_dark_mode,
                                         bootstyle="round-toggle")
        dark_mode_check.pack(side=LEFT, padx=(0, 15))
        
        ttk.Label(theme_frame, text="Theme:", font=('Segoe UI', 10)).pack(side=LEFT, padx=(0, 5))
        
        self.theme_dropdown = ttk.Combobox(theme_frame, 
                                     textvariable=self.current_theme,
                                     values=self.themes,
                                     state="readonly",
                                     width=15)
        self.theme_dropdown.pack(side=LEFT)
        self.theme_dropdown.bind('<<ComboboxSelected>>', self.change_theme)
        
        # Main container with padding
        main_container = ttk.Frame(self.root, padding="20")
        main_container.pack(fill=BOTH, expand=TRUE)
        
        # Create paned window for resizable layout
        paned = ttk.PanedWindow(main_container, orient=HORIZONTAL)
        paned.pack(fill=BOTH, expand=TRUE)
        
        # Left panel (controls)
        left_panel = ttk.Frame(paned, padding="10")
        paned.add(left_panel, weight=1)
        
        # Right panel (preview) - add padding to create space from left panel
        right_panel = ttk.Frame(paned, padding=(20, 10, 10, 10))  # More padding on left side
        paned.add(right_panel, weight=2)
        
        # Create notebook for organized controls
        self.notebook = ttk.Notebook(left_panel, bootstyle="primary")
        self.notebook.pack(fill=BOTH, expand=TRUE)
        
        # Set initial sash position (after adding both panels)
        self.root.after(100, lambda: paned.sashpos(0, 520))  # Set left panel to 520px width
        
        # Create tabs
        self.create_input_tab()
        self.create_filter_tab()
        self.create_style_tab()
        
        # Create preview area
        self.create_preview_area(right_panel)
        
    def create_input_tab(self):
        """Create input sources tab"""
        input_frame = ttk.Frame(self.notebook, padding="20")
        self.notebook.add(input_frame, text="📁 Input")
        
        # Working folder selection
        folder_frame = self.create_section(input_frame, "Working Folder")
        
        folder_info = ttk.Frame(folder_frame)
        folder_info.pack(fill=X, pady=(0, 10))
        
        ttk.Label(folder_info, 
                 textvariable=self.working_folder,
                 bootstyle="secondary",
                 font=('Segoe UI', 10)).pack(side=LEFT, padx=(5, 0))
        
        ttk.Button(folder_frame, 
                  text="Select Folder",
                  command=self.select_folder,
                  bootstyle="primary",
                  width=20).pack()
        
        # File selection
        file_frame = self.create_section(input_frame, "Select Files")
        
        # Create frame for listbox with border
        listbox_frame = ttk.Frame(file_frame, bootstyle="secondary", padding=1)
        listbox_frame.pack(fill=BOTH, expand=TRUE, pady=(0, 10))
        
        self.file_listbox = tk.Listbox(listbox_frame,
                                      selectmode=tk.MULTIPLE,
                                      height=6,
                                      font=('Segoe UI', 10),
                                      borderwidth=0,
                                      highlightthickness=1,
                                      highlightbackground="#e0e0e0",
                                      highlightcolor="#0078d4")
        self.file_listbox.pack(fill=BOTH, expand=TRUE, padx=1, pady=1)
        
        # Button frame for file operations
        file_btn_frame = ttk.Frame(file_frame)
        file_btn_frame.pack(fill=X, pady=(5, 0))
        
        ttk.Button(file_btn_frame,
                  text="Select All",
                  command=self.select_all_files,
                  bootstyle="info-outline",
                  width=10).pack(side=LEFT, padx=(0, 5))
        
        ttk.Button(file_btn_frame,
                  text="Clear",
                  command=self.clear_file_selection,
                  bootstyle="secondary-outline",
                  width=12).pack(side=LEFT, padx=(0, 5))
        
        ttk.Button(file_btn_frame,
                  text="Load Selected Files",
                  command=self.load_files,
                  bootstyle="success",
                  width=18).pack(side=LEFT)
        
        # Text input
        text_frame = self.create_section(input_frame, "Or Paste Text")
        
        # Create frame for text widget with border
        text_border = ttk.Frame(text_frame, bootstyle="secondary", padding=1)
        text_border.pack(fill=BOTH, expand=TRUE, pady=(0, 10))
        
        self.text_input = ScrolledText(text_border,
                                      height=8,
                                      font=('Segoe UI', 10),
                                      borderwidth=0,
                                      highlightthickness=0,
                                      wrap=tk.WORD)
        self.text_input.pack(fill=BOTH, expand=TRUE, padx=1, pady=1)
        
        ttk.Button(text_frame,
                  text="Use Pasted Text",
                  command=self.use_pasted_text,
                  bootstyle="info",
                  width=20).pack()
        
    def create_filter_tab(self):
        """Create filters tab"""
        filter_frame = ttk.Frame(self.notebook, padding="20")
        self.notebook.add(filter_frame, text="⚙️ Filters")
        
        # Word length filters
        length_frame = self.create_section(filter_frame, "Word Length")
        
        # Create a horizontal frame for both meters
        meters_frame = ttk.Frame(length_frame)
        meters_frame.pack(fill=X, pady=(0, 20))
        
        # Min length meter
        min_container = ttk.Frame(meters_frame)
        min_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 20))
        
        try:
            # Add label above meter
            ttk.Label(min_container, text="Minimum Length", 
                     font=('Segoe UI', 11, 'bold')).pack(pady=(0, 10))
            
            self.min_length_meter = Meter(
                min_container,
                metersize=180,
                amountused=3,
                amounttotal=50,
                metertype='semi',
                textleft='',
                textright='chars',
                interactive=True,
                bootstyle='primary',
                stripethickness=2  # Smooth continuous line
            )
            self.min_length_meter.pack()
            
            # Add description
            ttk.Label(min_container, text="Filter out words shorter than this", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(min_container, text="1 - 50", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Bind the meter value change
            self.min_length_meter.amountusedvar.trace('w', lambda *args: self.update_min_from_meter())
            self.min_length_scale = None  # Not using scale
        except Exception as e:
            # Fallback to scale if meter fails
            debug_print(f"Meter creation failed: {e}, using scale instead")
            min_label_frame = ttk.Frame(min_container)
            min_label_frame.pack(fill=X)
            ttk.Label(min_label_frame, text="Minimum Length:", 
                     font=('Segoe UI', 11, 'bold')).pack(side=LEFT)
            self.min_length_label = ttk.Label(min_label_frame, text="3",
                                             bootstyle="primary", 
                                             font=('Segoe UI', 14, 'bold'))
            self.min_length_label.pack(side=RIGHT)
            
            self.min_length_scale = ttk.Scale(min_container,
                                             from_=1,
                                             to=50,
                                             value=3,
                                             command=self.update_min_label,
                                             bootstyle="primary",
                                             length=200)
            self.min_length_scale.pack(fill=X, pady=(10, 0))
            self.min_length_meter = None
        
        # Max length meter
        max_container = ttk.Frame(meters_frame)
        max_container.pack(side=LEFT, fill=BOTH, expand=True)
        
        try:
            # Add label above meter
            ttk.Label(max_container, text="Maximum Length", 
                     font=('Segoe UI', 11, 'bold')).pack(pady=(0, 10))
            
            self.max_length_meter = Meter(
                max_container,
                metersize=180,
                amountused=30,
                amounttotal=50,
                metertype='semi',
                textleft='',
                textright='chars',
                interactive=True,
                bootstyle='info',
                stripethickness=2  # Smooth continuous line
            )
            self.max_length_meter.pack()
            
            # Add description
            ttk.Label(max_container, text="Filter out words longer than this", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(max_container, text="1 - 50", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Bind the meter value change
            self.max_length_meter.amountusedvar.trace('w', lambda *args: self.update_max_from_meter())
            self.max_length_scale = None  # Not using scale
        except Exception as e:
            # Fallback to scale if meter fails
            debug_print(f"Meter creation failed: {e}, using scale instead")
            max_label_frame = ttk.Frame(max_container)
            max_label_frame.pack(fill=X)
            ttk.Label(max_label_frame, text="Maximum Length:", 
                     font=('Segoe UI', 11, 'bold')).pack(side=LEFT)
            self.max_length_label = ttk.Label(max_label_frame, text="30",
                                             bootstyle="info", 
                                             font=('Segoe UI', 14, 'bold'))
            self.max_length_label.pack(side=RIGHT)
            
            self.max_length_scale = ttk.Scale(max_container,
                                             from_=3,
                                             to=50,
                                             value=30,
                                             command=self.update_max_label,
                                             bootstyle="info",
                                             length=200)
            self.max_length_scale.pack(fill=X, pady=(10, 0))
            self.max_length_meter = None
        
        # Forbidden words
        forbidden_frame = self.create_section(filter_frame, "Forbidden Words")
        
        ttk.Label(forbidden_frame,
                 text="Enter words to exclude (one per line):",
                 font=('Segoe UI', 10)).pack(anchor=W, pady=(0, 5))
        
        # Create frame for text widget with border
        text_border = ttk.Frame(forbidden_frame, bootstyle="secondary", padding=1)
        text_border.pack(fill=BOTH, expand=TRUE, pady=(0, 10))
        
        self.forbidden_text = ScrolledText(text_border,
                                          height=10,
                                          font=('Segoe UI', 10),
                                          borderwidth=0,
                                          highlightthickness=0,
                                          wrap=tk.WORD)
        self.forbidden_text.pack(fill=BOTH, expand=TRUE, padx=1, pady=1)
        
        # Store default forbidden words for reset - using wordcloud STOPWORDS
        self.default_forbidden = '\n'.join(sorted(STOPWORDS))
        
        # Don't pre-populate here - let config loading handle it
        # If no config is loaded, we'll insert defaults later
        
        ttk.Button(forbidden_frame,
                  text="Update Forbidden Words",
                  command=self.update_forbidden_words,
                  bootstyle="warning",
                  width=25).pack()
        
    def create_style_tab(self):
        """Create style options tab"""
        style_tab = ttk.Frame(self.notebook)
        self.notebook.add(style_tab, text="🎨 Style")
        
        # Create scrollable frame
        canvas = tk.Canvas(style_tab, highlightthickness=0)
        scrollbar = ttk.Scrollbar(style_tab, orient="vertical", command=canvas.yview, bootstyle="primary-round")
        scrollable_frame = ttk.Frame(canvas)
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        # Create the window and store its ID
        self.style_window_id = canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # Function to update canvas window width
        def update_style_canvas_width(event=None):
            canvas_width = canvas.winfo_width()
            if canvas_width > 1:  # Ensure canvas has been drawn
                canvas.itemconfig(self.style_window_id, width=canvas_width)
        
        # Bind canvas resize to update window width
        canvas.bind("<Configure>", update_style_canvas_width)
        
        # Pack scrollbar and canvas
        scrollbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)
        
        # Update width after canvas is displayed
        canvas.after(100, update_style_canvas_width)
        
        # Add padding to scrollable frame
        style_frame = ttk.Frame(scrollable_frame, padding="20")
        style_frame.pack(fill="both", expand=True)
        
        # Bind mouse wheel to this specific canvas
        def _on_style_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        
        # Store the current binding
        self._style_wheel_bound = False
        
        def _bind_style_mousewheel(event):
            if not self._style_wheel_bound:
                canvas.bind("<MouseWheel>", _on_style_mousewheel)
                self._style_wheel_bound = True
        
        def _unbind_style_mousewheel(event):
            if self._style_wheel_bound:
                canvas.unbind("<MouseWheel>")
                self._style_wheel_bound = False
        
        # Bind/unbind mousewheel when entering/leaving the canvas
        canvas.bind('<Enter>', _bind_style_mousewheel)
        canvas.bind('<Leave>', _unbind_style_mousewheel)
        
        # Also bind to the scrollable frame and its children
        scrollable_frame.bind('<MouseWheel>', _on_style_mousewheel)
        style_frame.bind('<MouseWheel>', _on_style_mousewheel)
        
        # Color scheme selection
        color_frame = self.create_section(style_frame, "Color Scheme")
        
        
        # Create notebook for different color modes
        self.color_notebook = ttk.Notebook(color_frame, bootstyle="primary")
        self.color_notebook.pack(fill=BOTH, expand=TRUE)
        
        # Bind tab change event
        self.color_notebook.bind("<<NotebookTabChanged>>", self.on_color_tab_changed)
        
        # Single color tab
        single_tab = ttk.Frame(self.color_notebook)
        self.color_notebook.add(single_tab, text="Single Color")
        
        single_color_frame = ttk.Frame(single_tab, padding=20)
        single_color_frame.pack(fill=X)
        
        ttk.Label(single_color_frame, text="Color:", font=('Segoe UI', 10)).pack(side=LEFT)
        
        self.single_color_preview = ttk.Frame(single_color_frame, width=30, height=30)
        self.single_color_preview.pack(side=LEFT, padx=(10, 10))
        
        # Set initial color preview
        style = ttk.Style()
        style_name = "SingleColorPreview.TFrame"
        style.configure(style_name, background=self.single_color.get())
        self.single_color_preview.configure(style=style_name)
        
        self.single_color_btn = ttk.Button(single_color_frame,
                                         text="Choose Color",
                                         command=self.choose_single_color,
                                         bootstyle="primary-outline")
        self.single_color_btn.pack(side=LEFT)
        
        # Preset gradients tab
        preset_tab = ttk.Frame(self.color_notebook)
        self.color_notebook.add(preset_tab, text="Preset Gradients")
        
        # Create scrollable frame for preset color buttons
        preset_canvas = tk.Canvas(preset_tab, height=300)
        preset_scrollbar = ttk.Scrollbar(preset_tab, orient="vertical", command=preset_canvas.yview, bootstyle="primary-round")
        preset_scrollable = ttk.Frame(preset_canvas)
        
        preset_scrollable.bind(
            "<Configure>",
            lambda e: preset_canvas.configure(scrollregion=preset_canvas.bbox("all"))
        )
        
        # Create the window and store its ID
        self.preset_window_id = preset_canvas.create_window((0, 0), window=preset_scrollable, anchor="nw")
        preset_canvas.configure(yscrollcommand=preset_scrollbar.set)
        
        # Function to update canvas window width
        def update_preset_canvas_width(event=None):
            canvas_width = preset_canvas.winfo_width()
            if canvas_width > 1:  # Ensure canvas has been drawn
                preset_canvas.itemconfig(self.preset_window_id, width=canvas_width)
        
        # Bind canvas resize to update window width
        preset_canvas.bind("<Configure>", update_preset_canvas_width)
        
        preset_canvas.pack(side="left", fill="both", expand=True)
        preset_scrollbar.pack(side="right", fill="y")
        
        # Update width after canvas is displayed
        preset_canvas.after(100, update_preset_canvas_width)
        
        # Bind mouse wheel to preset canvas only
        def _on_preset_mousewheel(event):
            preset_canvas.yview_scroll(int(-1*(event.delta/120)), "units")
            # Stop propagation
            return "break"
        
        # Direct binding to the preset canvas
        preset_canvas.bind("<MouseWheel>", _on_preset_mousewheel)
        
        # Also bind to the scrollable frame inside
        preset_scrollable.bind("<MouseWheel>", _on_preset_mousewheel)
        
        # Custom gradient tab
        custom_tab = ttk.Frame(self.color_notebook)
        self.color_notebook.add(custom_tab, text="Custom Gradient")
        
        custom_frame = ttk.Frame(custom_tab, padding=20)
        custom_frame.pack(fill=BOTH, expand=TRUE)
        
        ttk.Label(custom_frame, text="Create your own gradient:", 
                 font=('Segoe UI', 10, 'bold')).pack(anchor=W, pady=(0, 10))
        
        # Custom gradient colors
        self.custom_color_frames = []
        self.custom_color_previews = []
        
        for i in range(3):
            color_row = ttk.Frame(custom_frame)
            color_row.pack(fill=X, pady=5)
            
            ttk.Label(color_row, text=f"Color {i+1}:", width=10).pack(side=LEFT)
            
            preview = ttk.Frame(color_row, width=30, height=30)
            preview.pack(side=LEFT, padx=(5, 10))
            self.custom_color_previews.append(preview)
            
            btn = ttk.Button(color_row, text="Choose", 
                           command=lambda idx=i: self.choose_custom_color(idx),
                           bootstyle="secondary-outline")
            btn.pack(side=LEFT)
            
            self.custom_color_frames.append(color_row)
        
        # Update custom color previews
        self.update_custom_gradient_preview()
        
        # Add/Remove color buttons
        btn_frame = ttk.Frame(custom_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        ttk.Button(btn_frame, text="Add Color", command=self.add_gradient_color,
                  bootstyle="success-outline").pack(side=LEFT, padx=(0, 10))
        
        ttk.Button(btn_frame, text="Remove Color", command=self.remove_gradient_color,
                  bootstyle="danger-outline").pack(side=LEFT)
        
        # Set initial tab based on color mode
        self.color_notebook.select(1)  # Select preset tab by default
        self.update_custom_gradient_preview()
        
        # Combined color preview frame (after the notebook)
        ttk.Separator(color_frame, orient='horizontal').pack(fill=X, pady=(10, 5))
        
        combined_preview_frame = ttk.Frame(color_frame)
        combined_preview_frame.pack(fill=X, pady=(5, 10))
        
        ttk.Label(combined_preview_frame, text="Selected Color Scheme Preview:", 
                 font=('Segoe UI', 10, 'bold')).pack(anchor=W, pady=(0, 5))
        
        self.combined_color_preview = ttk.Frame(combined_preview_frame, height=50)
        self.combined_color_preview.pack(fill=X)
        self.combined_color_preview.pack_propagate(False)
        # Create scrollable frame for color buttons
        color_scroll = preset_scrollable
        
        # Create color scheme buttons in a grid
        self.color_var = tk.StringVar(value="Viridis")
        
        colors_grid = ttk.Frame(color_scroll)
        colors_grid.pack(fill=X, padx=10, pady=(10, 0))
        
        # Bind mouse wheel to grid
        colors_grid.bind("<MouseWheel>", _on_preset_mousewheel)
        
        row = 0
        col = 0
        for name, cmap in self.color_schemes.items():
            btn = ttk.Radiobutton(colors_grid,
                                 text=name,
                                 variable=self.color_var,
                                 value=name,
                                 command=self.on_color_select,
                                 bootstyle="primary-outline-toolbutton",
                                 width=12)
            btn.grid(row=row, column=col, padx=5, pady=5, sticky=W)
            # Bind mouse wheel to button
            btn.bind("<MouseWheel>", _on_preset_mousewheel)
            col += 1
            if col > 3:  # Changed from 1 to 3 for 4 columns
                col = 0
                row += 1
        
        # Update combined preview after all color vars are initialized
        self.update_combined_color_preview()
        
        # Mask and Shape Options
        mask_frame = self.create_section(style_frame, "Shape & Appearance")
        # add a scroll bind
        mask_frame.bind("<MouseWheel>", _on_style_mousewheel)
        
        self.mask_type = tk.StringVar(value="no_mask")
        
        ttk.Separator(mask_frame, orient='horizontal').pack(fill=X, pady=(5, 10))
        
        # Create notebook for mask options
        self.mask_notebook = ttk.Notebook(mask_frame, bootstyle="primary")
        self.mask_notebook.pack(fill=BOTH, expand=TRUE)
        
        # Bind tab change event
        self.mask_notebook.bind("<<NotebookTabChanged>>", self.on_mask_tab_changed)
        
        # Create tabs
        self.create_no_mask_tab()
        self.create_image_mask_tab()
        self.create_text_mask_tab()
        
        # Bind tab change event
        self.mask_notebook.bind("<<NotebookTabChanged>>", self.on_mask_tab_changed)
        
        # Canvas options
        canvas_frame = ttk.LabelFrame(mask_frame, text="Canvas Settings", padding=10)
        canvas_frame.pack(fill=X, pady=(0, 10))
        
        
        # Center container
        canvas_center = ttk.Frame(canvas_frame)
        canvas_center.pack(expand=True)
        
        # Lock aspect ratio checkbox
        ratio_frame = ttk.Frame(canvas_center)
        ratio_frame.pack(pady=(0, 10))
        
        self.lock_ratio_check = ttk.Checkbutton(ratio_frame,
                                               text="Lock aspect ratio",
                                               variable=self.lock_aspect_ratio,
                                               command=self.on_lock_ratio_change,
                                               bootstyle="primary")
        self.lock_ratio_check.pack(side=LEFT)
        
        self.ratio_label = ttk.Label(ratio_frame, text="",
                                    font=('Segoe UI', 9, 'italic'),
                                    bootstyle="secondary")
        self.ratio_label.pack(side=LEFT, padx=(10, 0))
        
        # Create horizontal container for width and height meters
        dimensions_container = ttk.Frame(canvas_center)
        dimensions_container.pack(pady=(0, 15))
        
        # Width meter
        width_container = ttk.Frame(dimensions_container)
        width_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 10))
        
        
        try:
            ttk.Label(width_container, text="Width", 
                     font=('Segoe UI', 10, 'bold')).pack(pady=(0, 5))
            
            self.width_meter = Meter(
                width_container,
                metersize=120,
                amountused=800,
                amounttotal=4000,
                metertype='semi',
                textleft='',
                textright='px',
                interactive=True,
                bootstyle='primary',
                stripethickness=0  # Smooth continuous line
            )
            self.width_meter.pack()
            
            # Add description
            ttk.Label(width_container, text="Canvas width in pixels", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(width_container, text="100 - 4000", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.width_meter.amountusedvar.trace('w', lambda *args: self.update_width_from_meter())
            
            
            self.width_scale = None
            self.width_label = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Width meter failed: {e}")
            width_label_frame = ttk.Frame(width_container)
            width_label_frame.pack(fill=X)
            ttk.Label(width_label_frame, text="Width:", font=('Segoe UI', 10)).pack(side=LEFT)
            self.width_label = ttk.Label(width_label_frame, text="800 px",
                                        bootstyle="primary", font=('Segoe UI', 10, 'bold'))
            self.width_label.pack(side=RIGHT)
            
            self.width_scale = ttk.Scale(width_container,
                                        from_=400,
                                        to=4000,
                                        value=800,
                                        command=self.update_width,
                                        bootstyle="primary")
            self.width_scale.pack(fill=X, pady=(5, 0))
            self.width_meter = None
        
        # Height meter
        height_container = ttk.Frame(dimensions_container)
        height_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(10, 0))
        
        
        try:
            ttk.Label(height_container, text="Height", 
                     font=('Segoe UI', 10, 'bold')).pack(pady=(0, 5))
            
            self.height_meter = Meter(
                height_container,
                metersize=120,
                amountused=600,
                amounttotal=4000,
                metertype='semi',
                textleft='',
                textright='px',
                interactive=True,
                bootstyle='primary',
                stripethickness=0  # Smooth continuous line
            )
            self.height_meter.pack()
            
            # Add description
            ttk.Label(height_container, text="Canvas height in pixels", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(height_container, text="100 - 4000", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.height_meter.amountusedvar.trace('w', lambda *args: self.update_height_from_meter())
            
            
            self.height_scale = None
            self.height_label = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Height meter failed: {e}")
            height_label_frame = ttk.Frame(height_container)
            height_label_frame.pack(fill=X)
            ttk.Label(height_label_frame, text="Height:", font=('Segoe UI', 10)).pack(side=LEFT)
            self.height_label = ttk.Label(height_label_frame, text="600 px",
                                         bootstyle="primary", font=('Segoe UI', 10, 'bold'))
            self.height_label.pack(side=RIGHT)
            
            self.height_scale = ttk.Scale(height_container,
                                         from_=300,
                                         to=4000,
                                         value=600,
                                         command=self.update_height,
                                         bootstyle="primary")
            self.height_scale.pack(fill=X, pady=(5, 0))
            self.height_meter = None
        
        # Size presets
        preset_frame = ttk.Frame(canvas_center)
        preset_frame.pack(pady=(10, 20))
        
        ttk.Label(preset_frame, text="", font=('Segoe UI', 10)).pack(side=LEFT, padx=(0, 10))
        
        presets = [
            ("Square", 800, 800),
            ("HD", 1920, 1080),
            ("4:3", 800, 600),
            ("16:9", 1280, 720),
            ("A4", 800, 1131)
        ]
        
        for name, width, height in presets:
            ttk.Button(preset_frame,
                      text=name,
                      command=lambda w=width, h=height: self.set_canvas_size(w, h),
                      bootstyle="secondary-outline",
                      width=8).pack(side=LEFT, padx=2)
        
        # Mode and Background Color center container
        mode_bg_center = ttk.Frame(canvas_frame)
        mode_bg_center.pack(expand=True, pady=(10, 0))
        
        # Mode selection (RGB/RGBA)
        mode_container = ttk.Frame(mode_bg_center)
        mode_container.pack(fill=X, pady=(0, 10))
        
        ttk.Label(mode_container, text="Mode:", font=('Segoe UI', 10)).pack(side=LEFT)
        
        mode_frame = ttk.Frame(mode_container)
        mode_frame.pack(side=LEFT, padx=(20, 0))
        
        ttk.Radiobutton(mode_frame,
                       text="RGB (Solid)",
                       variable=self.rgba_mode,
                       value=False,
                       command=self.update_mode,
                       bootstyle="primary").pack(side=LEFT, padx=(0, 15))
        
        rgba_radio = ttk.Radiobutton(mode_frame,
                       text="RGBA (Transparent)",
                       variable=self.rgba_mode,
                       value=True,
                       command=self.update_mode,
                       bootstyle="primary")
        rgba_radio.pack(side=LEFT)
        
        # Show transparency checkbox (only visible in RGBA mode)
        self.transparency_check = ttk.Checkbutton(mode_frame,
                                                 text="Show transparency",
                                                 variable=self.show_transparency,
                                                 bootstyle="primary")
        self.transparency_check.pack(side=LEFT, padx=(10, 0))
        
        # Initialize checkbox state based on current mode
        if not self.rgba_mode.get():
            self.transparency_check.configure(state=DISABLED)
        
        # Background color
        self.bg_container = ttk.Frame(mode_bg_center)
        self.bg_container.pack(fill=X)
        
        self.bg_label = ttk.Label(self.bg_container, text="Background Color:", font=('Segoe UI', 10))
        self.bg_label.pack(side=LEFT)
        
        self.bg_color_preview = ttk.Frame(self.bg_container, width=30, height=30)
        self.bg_color_preview.pack(side=RIGHT, padx=(10, 0))
        self.bg_color_preview.configure(bootstyle="light")
        
        self.bg_color_btn = ttk.Button(self.bg_container,
                                      text="Choose Color",
                                      command=self.choose_bg_color,
                                      bootstyle="primary-outline",
                                      width=15)
        self.bg_color_btn.pack(side=RIGHT)
        
        # Word Orientation
        orientation_frame = ttk.LabelFrame(mask_frame, text="Word Orientation", padding=10)
        orientation_frame.pack(fill=X, pady=(0, 10))
        
        # Center container for all content
        center_container = ttk.Frame(orientation_frame)
        center_container.pack(expand=True)
        
        # Prefer horizontal with Floodgauge
        horizontal_container = ttk.Frame(center_container)
        horizontal_container.pack()
        
        ttk.Label(horizontal_container, text="Word Orientation", 
                 font=('Segoe UI', 11, 'bold')).pack()
        
        ttk.Label(horizontal_container, 
                 text="Control the ratio of horizontal to vertical words in your cloud",
                 font=('Segoe UI', 9),
                 bootstyle="secondary").pack(pady=(5, 0))
        
        # Create a frame to hold the gauge and labels
        gauge_container = ttk.Frame(horizontal_container)
        gauge_container.pack(pady=10)
        
        gauge_frame = ttk.Frame(gauge_container)
        gauge_frame.pack()
        
        # Left label
        ttk.Label(gauge_frame, text="Vertical", 
                 font=('Segoe UI', 9)).pack(side=LEFT, padx=(0, 10))
        
        # Create Floodgauge for orientation
        self.horizontal_gauge = Floodgauge(
            gauge_frame,
            length=200,
            maximum=100,
            value=90,
            mask="{}%",
            font=('Segoe UI', 10, 'bold'),
            bootstyle="primary",
            mode='determinate',
            orient='horizontal'
        )
        self.horizontal_gauge.pack(side=LEFT)
        
        # Right label
        ttk.Label(gauge_frame, text="Horizontal", 
                 font=('Segoe UI', 9)).pack(side=LEFT, padx=(10, 0))
        
        # Add interactive control
        control_frame = ttk.Frame(horizontal_container)
        control_frame.pack()
        
        scale_container = ttk.Frame(control_frame)
        scale_container.pack()
        
        self.horizontal_scale = ttk.Scale(scale_container,
                                        from_=0,
                                        to=100,
                                        value=90,
                                        command=self.update_horizontal_gauge,
                                        bootstyle="primary",
                                        length=250)
        self.horizontal_scale.pack(pady=(5, 10))
        
        # Reset button
        ttk.Button(control_frame,
                  text="Reset to Default (90%)",
                  command=lambda: self.reset_orientation(),
                  bootstyle="secondary-outline").pack()
        
        # Other Settings
        other_frame = ttk.LabelFrame(mask_frame, text="Other Settings", padding=10)
        other_frame.pack(fill=X, pady=(0, 10))
        
        # Center container
        center_container = ttk.Frame(other_frame)
        center_container.pack(expand=True)
        
        # Create a grid layout for meters
        meters_grid = ttk.Frame(center_container)
        meters_grid.pack()
        
        # Letter thickness meter (moved to text mask tab)
        self.letter_thickness = tk.DoubleVar(value=1.0)
        # Letter spacing variable
        self.letter_spacing = tk.DoubleVar(value=0.0)
        
        # Max words meter
        max_words_container = ttk.Frame(meters_grid)
        max_words_container.grid(row=0, column=0, padx=10, pady=10)
        
        try:
            # Add label above meter
            ttk.Label(max_words_container, text="Max Words", 
                     font=('Segoe UI', 10, 'bold')).pack(pady=(0, 5))
            
            self.max_words_meter = Meter(
                max_words_container,
                metersize=150,
                amountused=200,
                amounttotal=500,
                metertype='semi',
                textleft='',
                textright='words',
                interactive=True,
                bootstyle='success',
                stripethickness=0  # Smooth continuous line
            )
            self.max_words_meter.pack()
            
            # Add description
            ttk.Label(max_words_container, text="Maximum words to display", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(max_words_container, text="10 - 500", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.max_words_meter.amountusedvar.trace('w', lambda *args: self.update_max_words_from_meter())
            self.max_words_scale = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Max words meter failed: {e}")
            self.max_words_label = ttk.Label(max_words_container, text="200",
                                            bootstyle="primary", font=('Segoe UI', 10, 'bold'))
            self.max_words_label.pack()
            self.max_words_scale = ttk.Scale(max_words_container,
                                            from_=10,
                                            to=500,
                                            value=200,
                                            command=self.update_max_words,
                                            bootstyle="primary")
            self.max_words_scale.pack(fill=X, pady=(5, 0))
            self.max_words_meter = None
        
        ttk.Label(max_words_container, 
                 text="More words = denser cloud, fewer words = cleaner look",
                 font=('Segoe UI', 9),
                 bootstyle="secondary").pack(pady=(5, 0))
        
        # Scale meter
        scale_container = ttk.Frame(meters_grid)
        scale_container.grid(row=0, column=1, padx=10, pady=10)
        
        try:
            # Add label above meter
            ttk.Label(scale_container, text="Computation Scale", 
                     font=('Segoe UI', 10, 'bold')).pack(pady=(0, 5))
            
            self.scale_meter = Meter(
                scale_container,
                metersize=150,
                amountused=1,
                amounttotal=10,
                metertype='semi',
                textleft='',
                textright='scale',
                interactive=True,
                bootstyle='warning',
                stripethickness=0  # Smooth continuous line
            )
            self.scale_meter.pack()
            
            # Add description
            ttk.Label(scale_container, text="Higher = better quality, slower", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(scale_container, text="1 - 10", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.scale_meter.amountusedvar.trace('w', lambda *args: self.update_scale_from_meter())
            self.scale_scale = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Scale meter failed: {e}")
            self.scale_label = ttk.Label(scale_container, text="1",
                                        bootstyle="primary", font=('Segoe UI', 10, 'bold'))
            self.scale_label.pack()
            self.scale_scale = ttk.Scale(scale_container,
                                        from_=1,
                                        to=10,
                                        value=1,
                                        command=self.update_scale,
                                        bootstyle="primary")
            self.scale_scale.pack(fill=X, pady=(5, 0))
            self.scale_meter = None
        
        ttk.Label(scale_container, 
                 text="Higher = faster generation but coarser word placement",
                 font=('Segoe UI', 9),
                 bootstyle="secondary").pack(pady=(5, 0))
        
        # Words per line meter
        words_container = ttk.Frame(meters_grid)
        words_container.grid(row=1, column=0, padx=10, pady=10)
        
        try:
            ttk.Label(words_container, text="Words per Line", 
                     font=('Segoe UI', 10, 'bold')).pack(pady=(0, 5))
            
            self.words_per_line_meter = Meter(
                words_container,
                metersize=150,
                amountused=1,
                amounttotal=10,
                metertype='semi',
                textleft='',
                textright='words',
                interactive=True,
                bootstyle='info',
                stripethickness=0  # Smooth continuous line
            )
            self.words_per_line_meter.pack()
            
            # Add description
            ttk.Label(words_container, text="Words per line in text masks", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(words_container, text="1 - 10", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.words_per_line_meter.amountusedvar.trace('w', lambda *args: self.update_words_per_line_from_meter())
            self.words_per_line_scale = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Words per line meter failed: {e}")
            words_label_frame = ttk.Frame(words_container)
            words_label_frame.pack(fill=X)
            ttk.Label(words_label_frame, text="Words per line:", font=('Segoe UI', 10)).pack(side=LEFT)
            self.words_per_line_label = ttk.Label(words_label_frame, text="1 word",
                                                 bootstyle="primary", font=('Segoe UI', 10, 'bold'))
            self.words_per_line_label.pack(side=RIGHT)
            
            self.words_per_line_scale = ttk.Scale(words_container,
                                                  from_=1,
                                                  to=10,
                                                  value=1,
                                                  command=self.update_words_per_line,
                                                  bootstyle="primary")
            self.words_per_line_scale.pack(fill=X, pady=(5, 0))
            self.words_per_line_meter = None
        
        ttk.Label(words_container, 
                 text="Words per line for text masks",
                 font=('Segoe UI', 9),
                 bootstyle="secondary").pack(pady=(5, 0))
        
    def create_no_mask_tab(self):
        """Create the no mask tab"""
        no_mask_frame = ttk.Frame(self.mask_notebook)
        self.mask_notebook.add(no_mask_frame, text="No Mask")

        # Info frame with border
        info_frame = ttk.LabelFrame(no_mask_frame, text="Information", padding=15)
        info_frame.pack(fill=X, padx=20, pady=20)
        
        # Info label
        info_label = ttk.Label(info_frame, 
                              text="Word cloud will be generated in a rectangular shape.\nNo special shape or outlines will be applied.",
                              font=('Segoe UI', 10),
                              bootstyle="secondary")
        info_label.pack()
        
        # Add a note about using other tabs
        ttk.Label(info_frame,
                 text="\nTo use a custom shape, select the Image Mask or Text Mask tab.",
                 font=('Segoe UI', 9, 'italic'),
                 bootstyle="info").pack()
    
    def create_image_mask_tab(self):
        """Create the image mask tab"""
        image_mask_frame = ttk.Frame(self.mask_notebook, padding=20)
        self.mask_notebook.add(image_mask_frame, text="Image Mask")
        
        # Create the image mask frame content
        self.create_image_mask_frame(image_mask_frame)
        
        # Add outline options to this tab
        self.create_outline_options(image_mask_frame)
        
        # Add mask preview to this tab
        self.create_mask_preview(image_mask_frame)
    
    def create_text_mask_tab(self):
        """Create the text mask tab"""
        text_mask_frame = ttk.Frame(self.mask_notebook, padding=20)
        self.mask_notebook.add(text_mask_frame, text="Text Mask")
        
        # Create the text mask frame content
        self.create_text_mask_frame(text_mask_frame)
        
        # Add mask preview to this tab
        self.create_mask_preview(text_mask_frame)
    
    def create_outline_options(self, parent):
        """Create outline options frame"""
        # This function is now empty as outline options are moved to text_mask_frame
        pass
    
    def create_mask_preview(self, parent):
        """Create mask preview frame"""
        preview_container = ttk.LabelFrame(parent, text="Mask Preview", padding=10)
        preview_container.pack(fill=BOTH, expand=TRUE, pady=(10, 0))
        
        # Create a label for this specific tab
        preview_label = ttk.Label(preview_container,
                                 text="No mask selected",
                                 anchor=CENTER,
                                 font=('Segoe UI', 10))
        preview_label.pack(fill=BOTH, expand=TRUE)
        
        # Store reference based on parent tab
        if "image" in str(parent):
            self.image_mask_preview_label = preview_label
        else:
            self.text_mask_preview_label = preview_label
    
    def create_image_mask_frame(self, parent):
        """Create the image mask options frame"""
        mask_file_frame = ttk.LabelFrame(parent, text="Image File", padding=10)
        mask_file_frame.pack(fill=X)
        
        mask_info = ttk.Frame(mask_file_frame)
        mask_info.pack(fill=X, pady=(0, 10))
        
        self.image_mask_label = ttk.Label(mask_info,
                                         text="No image selected",
                                         bootstyle="secondary",
                                         font=('Segoe UI', 10))
        self.image_mask_label.pack(side=LEFT)
        
        mask_btn_frame = ttk.Frame(mask_file_frame)
        mask_btn_frame.pack(fill=X)
        
        ttk.Button(mask_btn_frame,
                  text="Select Image",
                  command=self.select_mask,
                  bootstyle="primary",
                  width=15).pack(side=LEFT, padx=(0, 10))
        
        ttk.Button(mask_btn_frame,
                  text="Clear",
                  command=self.clear_mask,
                  bootstyle="secondary",
                  width=15).pack(side=LEFT)
    
    def create_text_mask_frame(self, parent):
        """Create the text mask options frame"""
        text_input_frame = ttk.LabelFrame(parent, text="Text Input", padding=10)
        
        text_input_frame.pack(fill=X)
        
        # Text input
        ttk.Label(text_input_frame, text="Enter text for mask:", font=('Segoe UI', 10)).pack(anchor=W, pady=(0, 5))
        
        self.text_mask_entry = ttk.Entry(text_input_frame,
                                        textvariable=self.text_mask_input,
                                        font=('Segoe UI', 12),
                                        bootstyle="primary")
        self.text_mask_entry.pack(fill=X, pady=(0, 10))
        self.text_mask_entry.bind('<KeyRelease>', lambda e: self.update_text_mask())
        
        # Font selection
        font_frame = ttk.LabelFrame(text_input_frame, text="Font", padding=10)
        font_frame.pack(fill=X, pady=(0, 10))
        
        self.font_listbox = FontListbox(font_frame,
                                       self.available_fonts,
                                       textvariable=self.text_mask_font,
                                       width=35,
                                       height=5)
        self.font_listbox.pack(fill=X)
        self.font_listbox.bind('<<FontSelected>>', lambda e: self.update_text_mask())
        
        # Create horizontal container for font settings and outline settings
        meters_container = ttk.Frame(text_input_frame)
        meters_container.pack(fill=X, pady=(15, 0))
        
        # Font Settings frame (left side)
        font_settings_frame = ttk.LabelFrame(meters_container, text="Font Settings", padding=10)
        font_settings_frame.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 20))
        
        # Font size meter in font settings frame
        font_size_container = ttk.Frame(font_settings_frame)
        font_size_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 15))
        
        try:
            ttk.Label(font_size_container, text="Size", 
                     font=('Segoe UI', 9, 'bold')).pack(pady=(0, 5))
            
            self.font_size_meter = Meter(
                font_size_container,
                metersize=100,
                amountused=200,
                amounttotal=2000,
                metertype='semi',
                textleft='',
                textright='pt',
                interactive=True,
                bootstyle='info',
                stripethickness=0  # Smooth continuous line
            )
            self.font_size_meter.pack()
            
            # Add description
            ttk.Label(font_size_container, text="Font size for text mask", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(font_size_container, text="10 - 2000", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.font_size_meter.amountusedvar.trace('w', lambda *args: self.update_font_size_from_meter())
            self.font_size_scale = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Font size meter failed: {e}")
            font_size_label_frame = ttk.Frame(font_size_container)
            font_size_label_frame.pack(fill=X)
            ttk.Label(font_size_label_frame, text="Size:", font=('Segoe UI', 9)).pack(side=LEFT)
            self.font_size_label = ttk.Label(font_size_label_frame, text="200",
                                            bootstyle="primary", font=('Segoe UI', 9, 'bold'))
            self.font_size_label.pack(side=RIGHT)
            
            self.font_size_scale = ttk.Scale(font_size_container,
                                            from_=50,
                                            to=2000,
                                            value=200,
                                            command=self.update_font_size,
                                            bootstyle="primary")
            self.font_size_scale.pack(fill=X, pady=(5, 0))
            self.font_size_meter = None
        
        # Font style options in font settings frame
        style_container = ttk.Frame(font_settings_frame)
        style_container.pack(side=LEFT, fill=X, padx=(0, 10))
        
        ttk.Label(style_container, text="Style", 
                 font=('Segoe UI', 9, 'bold')).pack(pady=(0, 10))
        
        ttk.Checkbutton(style_container,
                       text="Bold",
                       variable=self.text_mask_bold,
                       command=self.update_text_mask,
                       bootstyle="primary").pack(anchor=W, pady=(0, 5))
        
        ttk.Checkbutton(style_container,
                       text="Italic",
                       variable=self.text_mask_italic,
                       command=self.update_text_mask,
                       bootstyle="primary").pack(anchor=W)
        
        # Text Appearance frame (middle)
        text_appearance_frame = ttk.LabelFrame(meters_container, text="Text Appearance", padding=10)
        text_appearance_frame.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 10))
        
        # Letter thickness meter
        thickness_container = ttk.Frame(text_appearance_frame)
        thickness_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 10))
        
        try:
            ttk.Label(thickness_container, text="Thickness", 
                     font=('Segoe UI', 9, 'bold')).pack(pady=(0, 5))
            
            self.thickness_meter = Meter(
                thickness_container,
                metersize=100,
                amountused=1,
                amounttotal=5,
                metertype='semi',
                textleft='',
                textright='px',
                interactive=True,
                bootstyle='warning',
                stripethickness=0  # Smooth continuous line
            )
            self.thickness_meter.pack()
            
            # Add description
            ttk.Label(thickness_container, text="Letter stroke thickness", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(thickness_container, text="0 - 5", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.thickness_meter.amountusedvar.trace('w', lambda *args: self.update_thickness_from_meter())
            
            # Add zero button for thickness
            def set_thickness_zero():
                # Directly set to 0
                if hasattr(self, 'thickness_meter') and self.thickness_meter:
                    self.thickness_meter.configure(amountused=0)
                self.letter_thickness.set(0)
                self.update_text_mask()
            
            zero_btn = ttk.Button(thickness_container, text="Reset", width=6, 
                                 command=set_thickness_zero,
                                 bootstyle="secondary-outline")
            zero_btn.pack(pady=(5, 0))
        except Exception as e:
            print(f"[ERROR] Thickness meter creation failed: {e}")
            debug_print(f"Thickness meter creation failed: {e}")
            self.thickness_meter = None
        
        # Letter spacing meter
        spacing_container = ttk.Frame(text_appearance_frame)
        spacing_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(10, 0))
        
        try:
            ttk.Label(spacing_container, text="Spacing", 
                     font=('Segoe UI', 9, 'bold')).pack(pady=(0, 5))
            
            self.spacing_meter = Meter(
                spacing_container,
                metersize=100,
                amountused=0,
                amounttotal=5,
                metertype='semi',
                textleft='',
                textright='px',
                interactive=True,
                bootstyle='info',
                stripethickness=0  # Smooth continuous line
            )
            self.spacing_meter.pack()
            
            # Add description
            ttk.Label(spacing_container, text="Space between letters", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(spacing_container, text="0 - 5", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.spacing_meter.amountusedvar.trace('w', lambda *args: self.update_spacing_from_meter())
            
            # Add zero button for spacing
            def set_spacing_zero():
                # Directly set to 0
                self.spacing_meter.configure(amountused=0)
                self.letter_spacing.set(0)
            
            zero_btn = ttk.Button(spacing_container, text="Reset", width=6, 
                                 command=set_spacing_zero,
                                 bootstyle="secondary-outline")
            zero_btn.pack(pady=(5, 0))
        except Exception as e:
            debug_print(f"Spacing meter creation failed: {e}")
            self.spacing_meter = None
        
        # Outline settings frame (right side)
        outline_frame = ttk.LabelFrame(meters_container, text="Outline Settings", padding=10)
        outline_frame.pack(side=LEFT, fill=BOTH, expand=True)
        
        # Create horizontal layout inside outline frame
        outline_layout = ttk.Frame(outline_frame)
        outline_layout.pack(fill=X)
        
        # Outline width meter
        width_container = ttk.Frame(outline_layout)
        width_container.pack(side=LEFT, fill=BOTH, expand=True, padx=(0, 15))
        
        try:
            ttk.Label(width_container, text="Width", 
                     font=('Segoe UI', 9, 'bold')).pack(pady=(0, 5))
            
            self.outline_width_meter = Meter(
                width_container,
                metersize=100,
                amountused=0.0,
                amounttotal=30,
                metertype='semi',
                textleft='',
                textright='px',
                interactive=True,
                bootstyle='primary',
                stripethickness=0,  # Smooth continuous line
                stepsize=1
            )
            self.outline_width_meter.pack()
            
            # Add description
            ttk.Label(width_container, text="Outline thickness around shape", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            
            # Add min/max values below meter
            ttk.Label(width_container, text="0 - 30", 
                     font=('Segoe UI', 8), foreground='gray').pack(pady=(5, 0))
            self.outline_width_meter.amountusedvar.trace('w', lambda *args: self.update_outline_width_from_meter())
            
            # Add zero button function
            def set_outline_zero():
                # Directly set to 0
                self.outline_width_meter.configure(amountused=0)
                self.outline_width.set(0)
            

            zero_btn = ttk.Button(width_container, text="Reset",
                                command=set_outline_zero,
                                bootstyle="secondary-outline")
            zero_btn.pack(pady=(5, 0))
            self.outline_width_scale = None
            self.outline_width_label = None
        except Exception as e:
            # Fallback to scale
            debug_print(f"Outline width meter failed: {e}")
            width_label_frame = ttk.Frame(width_container)
            width_label_frame.pack(fill=X)
            outline_width_lbl = ttk.Label(width_label_frame, text="Width:", font=('Segoe UI', 9))
            outline_width_lbl.pack(side=LEFT)
            outline_width_label = ttk.Label(width_label_frame, text="0 px",
                                           bootstyle="primary", font=('Segoe UI', 9, 'bold'))
            outline_width_label.pack(side=RIGHT)
            
            outline_width_scale = ttk.Scale(width_container,
                                           from_=0,
                                           to=30,
                                           value=0,
                                           command=lambda v: self.update_outline_width(v, outline_width_label),
                                           bootstyle="primary")
            outline_width_scale.pack(fill=X, pady=(5, 0))
            self.outline_width_label = outline_width_label
            self.outline_width_scale = outline_width_scale
            self.outline_width_meter = None
        
        # Outline color
        color_container = ttk.Frame(outline_layout)
        color_container.pack(side=LEFT, fill=X)
        
        ttk.Label(color_container, text="Color", 
                 font=('Segoe UI', 9, 'bold')).pack(pady=(0, 5))
        
        color_frame = ttk.Frame(color_container)
        color_frame.pack()
        
        outline_color_preview = ttk.Frame(color_frame, width=25, height=25, bootstyle="dark")
        outline_color_preview.pack(side=LEFT, padx=(0, 8))
        
        outline_color_btn = ttk.Button(color_frame,
                                      text="Choose",
                                      command=lambda: self.choose_outline_color(outline_color_preview),
                                      bootstyle="primary-outline",
                                      width=10)
        outline_color_btn.pack(side=LEFT)
        
        # Store reference
        self.outline_color_preview = outline_color_preview
    
    def on_mask_type_change(self):
        """Handle mask type radio button change"""
        # Don't clear canvas when mask type changes
        # self.clear_canvas()
        
        # Update the mode label to reflect the mask selection
        self.update_mode_label()
        
        # Update outline state based on new selection
        self.update_outline_state()
    
    def on_mask_tab_changed(self, event):
        """Handle mask tab change"""
        # Tab changes don't affect the radio button selection
        # This allows users to explore different mask options without committing
        pass
    
    
    def update_font_size(self, value):
        """Update font size label and regenerate text mask"""
        val = int(float(value))
        self.text_mask_font_size.set(val)
        if hasattr(self, 'font_size_label'):
            self.font_size_label.config(text=str(val))
        # Update text mask if we're on the text mask tab and have text
        if self.text_mask_input.get():
            self.update_text_mask()
    
    def update_words_per_line(self, value):
        """Update words per line label and regenerate text mask"""
        val = int(float(value))
        self.text_mask_words_per_line.set(val)
        if hasattr(self, 'words_per_line_label'):
            if val == 1:
                self.words_per_line_label.config(text="1 word")
            else:
                self.words_per_line_label.config(text=f"{val} words")
        if self.text_mask_input.get():
            self.update_text_mask()
    
    def get_ratio_text(self, width, height):
        """Get a readable aspect ratio text"""
        # Calculate GCD to simplify ratio
        from math import gcd
        g = gcd(width, height)
        w = width // g
        h = height // g
        
        # Check for common ratios
        common_ratios = {
            (16, 9): "16:9",
            (9, 16): "9:16",
            (4, 3): "4:3",
            (3, 4): "3:4",
            (16, 10): "16:10",
            (1, 1): "1:1",
            (3, 2): "3:2",
            (2, 3): "2:3"
        }
        
        if (w, h) in common_ratios:
            return common_ratios[(w, h)]
        
        # Simplify further if numbers are too large
        while w > 20 or h > 20:
            if w % 2 == 0 and h % 2 == 0:
                w //= 2
                h //= 2
            else:
                break
        
        return f"{w}:{h}"
    
    def on_lock_ratio_change(self):
        """Handle lock aspect ratio checkbox change"""
        if self.lock_aspect_ratio.get():
            # Calculate and store current aspect ratio
            width = self.canvas_width.get()
            height = self.canvas_height.get()
            if height > 0:
                self.aspect_ratio = width / height
                # Show ratio in simplified form
                ratio_text = self.get_ratio_text(width, height)
                self.ratio_label.config(text=f"({ratio_text})")
        else:
            self.ratio_label.config(text="")
    
    def update_width(self, value):
        """Update width and maintain aspect ratio if locked"""
        if hasattr(self, '_updating'):  # Prevent recursion
            return
            
        val = int(float(value))
        self.canvas_width.set(val)
        self.width_label.config(text=f"{val} px")
        
        if self.lock_aspect_ratio.get() and self.aspect_ratio > 0:
            self._updating = True
            try:
                # Update height to maintain aspect ratio
                new_height = int(val / self.aspect_ratio)
                new_height = max(300, min(4000, new_height))  # Clamp to valid range
                self.canvas_height.set(new_height)
                self.height_label.config(text=f"{new_height} px")
                self.height_scale.set(new_height)
            finally:
                delattr(self, '_updating')
        
        # Clear canvas when dimensions change
        self.clear_canvas()
    
    def update_height(self, value):
        """Update height and maintain aspect ratio if locked"""
        if hasattr(self, '_updating'):  # Prevent recursion
            return
            
        val = int(float(value))
        self.canvas_height.set(val)
        self.height_label.config(text=f"{val} px")
        
        if self.lock_aspect_ratio.get() and self.aspect_ratio > 0:
            self._updating = True
            try:
                # Update width to maintain aspect ratio
                new_width = int(val * self.aspect_ratio)
                new_width = max(400, min(4000, new_width))  # Clamp to valid range
                self.canvas_width.set(new_width)
                self.width_label.config(text=f"{new_width} px")
                self.width_scale.set(new_width)
            finally:
                delattr(self, '_updating')
        
        # Clear canvas when dimensions change
        self.clear_canvas()
    
    def set_canvas_size(self, width, height):
        """Set canvas size from preset"""
        # Update the aspect ratio if locked
        if self.lock_aspect_ratio.get():
            self.aspect_ratio = width / height
            # Update ratio display
            ratio_text = self.get_ratio_text(width, height)
            self.ratio_label.config(text=f"({ratio_text})")
        
        # Update values and UI
        self.canvas_width.set(width)
        self.canvas_height.set(height)
        
        # Update meters or scales
        if self.width_meter:
            self.width_meter.amountusedvar.set(width)
        elif self.width_scale:
            self.width_label.config(text=f"{width} px")
            self.width_scale.set(width)
            
        if self.height_meter:
            self.height_meter.amountusedvar.set(height)
        elif self.height_scale:
            self.height_label.config(text=f"{height} px")
            self.height_scale.set(height)
        
        # Show toast with preset info
        ratio_text = self.get_ratio_text(width, height)
        self.show_toast(f"Canvas size set to {width}×{height} ({ratio_text})", "info")
        
        # Clear canvas when preset is selected (canvas size change)
        self.clear_canvas()
        
    def calculate_preview_size(self):
        """Calculate preview display size maintaining aspect ratio with max constraints"""
        actual_width = self.canvas_width.get()
        actual_height = self.canvas_height.get()
        
        # Apply user's preview scale preference
        user_scale = self.preview_scale.get() / 100.0
        scaled_width = actual_width * user_scale
        scaled_height = actual_height * user_scale
        
        # Define max constraints for preview
        max_preview_width = 600
        max_preview_height = 450  # Reasonable height limit
        
        # Calculate scale factors for both dimensions after user scaling
        width_scale = max_preview_width / scaled_width if scaled_width > max_preview_width else 1.0
        height_scale = max_preview_height / scaled_height if scaled_height > max_preview_height else 1.0
        
        # Use the smaller scale to ensure it fits in both dimensions
        constraint_scale = min(width_scale, height_scale)
        
        # Final scale is user scale * constraint scale
        final_scale = user_scale * constraint_scale
        
        display_width = int(actual_width * final_scale)
        display_height = int(actual_height * final_scale)
        
        self.print_debug(f"Preview size: {display_width}x{display_height} (user scale: {user_scale:.2f}, constraint scale: {constraint_scale:.2f}, final: {final_scale:.2f})")
        
        return display_width, display_height
    
    def create_preview_area(self, parent):
        """Create the word cloud preview area"""
        preview_container = ttk.LabelFrame(parent, text="Word Cloud Preview", padding=15)
        preview_container.pack(fill=BOTH, expand=TRUE)
        
        # Create a centered frame for the preview with margins
        preview_wrapper = ttk.Frame(preview_container)
        preview_wrapper.pack(fill=BOTH, expand=TRUE, padx=10)  # Reduced horizontal margins
        
        # Modern status bar header
        header_container = ttk.Frame(preview_wrapper)
        header_container.pack(fill=X, pady=(0, 12))
        
        # Create a custom styled frame for the status bar
        self.status_bar = tk.Frame(header_container, bg='#F3F4F6', height=44)
        self.status_bar.pack(fill=X)
        self.status_bar.pack_propagate(False)
        status_bar = self.status_bar  # Keep local reference for convenience
        
        # Add subtle top border
        self.status_bar_top_border = tk.Frame(status_bar, bg='#E5E7EB', height=1)
        self.status_bar_top_border.pack(fill=X, side=TOP)
        
        # Inner container with padding
        self.status_bar_inner = tk.Frame(status_bar, bg='#F3F4F6')
        self.status_bar_inner.pack(fill=BOTH, expand=True, padx=20)
        inner_container = self.status_bar_inner
        
        # Left side - Source info
        left_frame = tk.Frame(inner_container, bg='#F3F4F6')
        left_frame.pack(side=LEFT, fill=Y)
        
        # Source container
        source_container = tk.Frame(left_frame, bg='#F3F4F6')
        source_container.pack(expand=True)
        
        # Source icon and label
        source_row = tk.Frame(source_container, bg='#F3F4F6')
        source_row.pack()
        
        tk.Label(source_row, text="📁", font=('Segoe UI', 14), bg='#F3F4F6', fg='#6B7280').pack(side=LEFT, padx=(0, 8))
        
        source_text_frame = tk.Frame(source_row, bg='#F3F4F6')
        source_text_frame.pack(side=LEFT)
        
        tk.Label(source_text_frame, text="SOURCE", font=('Segoe UI', 8), bg='#F3F4F6', fg='#9CA3AF').pack(anchor='w')
        self.source_label = tk.Label(source_text_frame, text="None", font=('Segoe UI', 11, 'bold'), bg='#F3F4F6', fg='#1F2937')
        self.source_label.pack(anchor='w')
        
        # Center divider
        divider = tk.Frame(inner_container, bg='#E5E7EB', width=1)
        divider.pack(side=LEFT, fill=Y, padx=30)
        
        # Right side - Mask info
        right_frame = tk.Frame(inner_container, bg='#F3F4F6')
        right_frame.pack(side=LEFT, fill=Y)
        
        # Mask container
        mask_container = tk.Frame(right_frame, bg='#F3F4F6')
        mask_container.pack(expand=True)
        
        # Mask icon and label
        mask_row = tk.Frame(mask_container, bg='#F3F4F6')
        mask_row.pack()
        
        tk.Label(mask_row, text="🎭", font=('Segoe UI', 14), bg='#F3F4F6', fg='#6B7280').pack(side=LEFT, padx=(0, 8))
        
        mask_text_frame = tk.Frame(mask_row, bg='#F3F4F6')
        mask_text_frame.pack(side=LEFT)
        
        tk.Label(mask_text_frame, text="MASK", font=('Segoe UI', 8), bg='#F3F4F6', fg='#9CA3AF').pack(anchor='w')
        self.mask_label = tk.Label(mask_text_frame, text="No Mask", font=('Segoe UI', 11, 'bold'), bg='#F3F4F6', fg='#1F2937')
        self.mask_label.pack(anchor='w')
        
        # Second divider
        divider2 = tk.Frame(inner_container, bg='#E5E7EB', width=1)
        divider2.pack(side=LEFT, fill=Y, padx=30)
        
        # Color scheme info
        color_frame = tk.Frame(inner_container, bg='#F3F4F6')
        color_frame.pack(side=LEFT, fill=Y)
        
        # Color container
        color_container = tk.Frame(color_frame, bg='#F3F4F6')
        color_container.pack(expand=True)
        
        # Color icon and label
        color_row = tk.Frame(color_container, bg='#F3F4F6')
        color_row.pack()
        
        tk.Label(color_row, text="🎨", font=('Segoe UI', 14), bg='#F3F4F6', fg='#6B7280').pack(side=LEFT, padx=(0, 8))
        
        color_text_frame = tk.Frame(color_row, bg='#F3F4F6')
        color_text_frame.pack(side=LEFT)
        
        tk.Label(color_text_frame, text="COLOR SCHEME", font=('Segoe UI', 8), bg='#F3F4F6', fg='#9CA3AF').pack(anchor='w')
        self.color_scheme_label = tk.Label(color_text_frame, text="Single Color", font=('Segoe UI', 11, 'bold'), bg='#F3F4F6', fg='#1F2937')
        self.color_scheme_label.pack(anchor='w')
        
        # Add bottom border
        bottom_border = tk.Frame(status_bar, bg='#E5E7EB', height=1)
        bottom_border.pack(fill=X, side=BOTTOM)
        
        # Scale indicator label (initially hidden)
        self.scale_indicator = ttk.Label(preview_wrapper, 
                                        text="",
                                        font=('Segoe UI', 9, 'italic'),
                                        bootstyle="secondary")
        self.scale_indicator.pack(pady=(0, 5))
        
        # Canvas for word cloud with max width constraint
        canvas_container = ttk.Frame(preview_wrapper)
        canvas_container.pack(expand=TRUE)  # Center it
        
        # Create border frame
        border_frame = ttk.Frame(canvas_container, bootstyle="secondary", padding=2)
        border_frame.pack(pady=(0, 15))
        
        # Canvas frame inside border
        canvas_frame = ttk.Frame(border_frame, bootstyle="light")
        canvas_frame.pack()
        
        # Calculate initial display size
        display_width, display_height = self.calculate_preview_size()
        
        self.figure = plt.Figure(figsize=(display_width/100, display_height/100), facecolor='white')
        self.canvas = FigureCanvasTkAgg(self.figure, master=canvas_frame)
        self.canvas_widget = self.canvas.get_tk_widget()
        self.canvas_widget.pack()  # Don't expand, keep fixed size
        
        # Initial empty plot with message
        ax = self.figure.add_subplot(111)
        ax.text(0.5, 0.5, '', 
                horizontalalignment='center', verticalalignment='center',
                transform=ax.transAxes, fontsize=14, color='gray')
        ax.axis('off')
        self.canvas.draw()
        
        # Store reference to preview canvas frame for theme updates
        self.preview_canvas_frame = canvas_frame
        self.preview_border_frame = border_frame
        
        # Preview size control
        size_control_frame = ttk.Frame(preview_wrapper)
        size_control_frame.pack(fill=X, pady=(10, 20))
        
        # Center the controls
        size_center = ttk.Frame(size_control_frame)
        size_center.pack()
        
        ttk.Label(size_center, text="Preview Size:", 
                 font=('Segoe UI', 10)).pack(side=LEFT, padx=(0, 10))
        
        # Size percentage variable is already initialized in __init__
        
        # Smaller size for zoom out
        ttk.Button(size_center, text="−", 
                  command=lambda: self.adjust_preview_size(-10),
                  bootstyle="secondary-outline",
                  width=3).pack(side=LEFT, padx=2)
        
        # Size slider
        self.preview_slider = ttk.Scale(size_center,
                                      from_=25,
                                      to=200,
                                      value=100,
                                      orient=HORIZONTAL,
                                      length=200,
                                      command=self.update_preview_size_from_slider,
                                      bootstyle="info")
        self.preview_slider.pack(side=LEFT, padx=10)
        
        # Size label
        self.preview_size_label = ttk.Label(size_center,
                                          text="100%",
                                          font=('Segoe UI', 10, 'bold'),
                                          bootstyle="info",
                                          width=5)
        self.preview_size_label.pack(side=LEFT, padx=(5, 10))
        
        # Larger size for zoom in
        ttk.Button(size_center, text="+", 
                  command=lambda: self.adjust_preview_size(10),
                  bootstyle="secondary-outline",
                  width=3).pack(side=LEFT, padx=2)
        
        # Reset button
        ttk.Button(size_center, text="Reset",
                  command=lambda: self.set_preview_size(100),
                  bootstyle="secondary-link").pack(side=LEFT, padx=(10, 0))
        
        # Calculate and set the actual initial preview scale
        actual_width = self.canvas_width.get()
        actual_height = self.canvas_height.get()
        if actual_width > 0:
            # Check if preview is constrained
            if display_width < actual_width or display_height < actual_height:
                actual_scale = int((display_width / actual_width) * 100)
                self.preview_scale.set(actual_scale)
                self.preview_size_label.config(text=f"{actual_scale}%")
                self.preview_slider.set(actual_scale)
                # Update scale indicator
                self.scale_indicator.config(text=f"Preview at {actual_scale}% (limited by screen size)")
        
        # Button frame (centered below preview)
        button_frame = ttk.Frame(preview_wrapper)
        button_frame.pack(fill=X)
        
        # Progress bar (initially hidden)
        self.progress = ttk.Progressbar(button_frame, 
                                       mode='indeterminate',
                                       bootstyle="success-striped",
                                       length=300)
        
        # Main buttons container with debug toggle on the right
        buttons_row = ttk.Frame(button_frame)
        buttons_row.pack(fill=X)
        
        # Generate and save buttons (left/center)
        btn_container = ttk.Frame(buttons_row)
        btn_container.pack(side=LEFT, expand=True)
        
        self.generate_btn = ttk.Button(btn_container,
                                      text="🚀 Generate Word Cloud",
                                      command=self.generate_wordcloud,
                                      bootstyle="success",
                                      width=25)
        self.generate_btn.pack(side=LEFT, padx=(0, 10))
        
        self.save_btn = ttk.Button(btn_container,
                                  text="💾 Save Image",
                                  command=self.save_wordcloud,
                                  bootstyle="primary",
                                  width=20,
                                  state=DISABLED)
        self.save_btn.pack(side=LEFT)
        
        # Clear button
        self.clear_btn = ttk.Button(btn_container,
                                  text="🗑️ Clear",
                                  command=self.clear_canvas,
                                  bootstyle="secondary",
                                  width=15)
        self.clear_btn.pack(side=LEFT, padx=(10, 0))
        
        # Debug toggle (bottom right)
        debug_container = ttk.Frame(buttons_row)
        debug_container.pack(side=RIGHT, padx=(0, 20))
        
        self.debug_var = tk.BooleanVar(value=self.debug_mode)
        self.debug_toggle = ttk.Checkbutton(debug_container,
                                           text="Debug Mode",
                                           variable=self.debug_var,
                                           command=self.toggle_debug_mode,
                                           bootstyle="primary-round-toggle")
        self.debug_toggle.pack()
    
    def create_message_bar(self, parent):
        """Create the message bar in the specified parent"""
        # Message bar frame
        self.message_frame = ttk.Frame(parent)
        self.message_frame.pack(side=LEFT, fill=X, expand=TRUE)
        
        # Message styles
        self.message_styles = {
            "good": {"icon": "✓", "bootstyle": "success", "bg": "#d4edda", "fg": "#155724", "border": "#c3e6cb"},
            "info": {"icon": "ℹ", "bootstyle": "info", "bg": "#d1ecf1", "fg": "#0c5460", "border": "#bee5eb"},
            "warning": {"icon": "⚠", "bootstyle": "warning", "bg": "#fff3cd", "fg": "#856404", "border": "#ffeaa7"},
            "fail": {"icon": "✗", "bootstyle": "danger", "bg": "#f8d7da", "fg": "#721c24", "border": "#f5c6cb"}
        }
        
        # Create message label (initially hidden)
        self.message_container = ttk.Frame(self.message_frame)
        
        self.message_icon_label = ttk.Label(self.message_container, font=('Segoe UI', 12, 'bold'))
        self.message_icon_label.pack(side=LEFT, padx=(10, 5))
        
        self.message_label = ttk.Label(self.message_container, font=('Segoe UI', 10))
        self.message_label.pack(side=LEFT, padx=(0, 10))
        
        # Close button
        self.message_close_btn = ttk.Button(self.message_container, 
                                           text="×",
                                           width=3,
                                           command=self.hide_message)
        self.message_close_btn.pack(side=RIGHT, padx=(0, 5))
        
        # Initially hide the message
        self.hide_message()
    
    def show_message(self, message, status="info"):
        """Show a message in the message bar"""
        if status not in self.message_styles:
            status = "info"
        
        # Print to console based on status
        if status in ["bad", "error"]:
            self.print_fail(message)
        elif status == "warning":
            self.print_warning(message)
        else:
            self.print_info(message)
        
        style = self.message_styles[status]
        
        # Update message content
        self.message_icon_label.config(text=style["icon"])
        self.message_label.config(text=message)
        
        # Apply styling based on theme
        if self.current_theme.get() in ["darkly", "superhero", "solar", "cyborg", "vapor"]:
            # Dark theme adjustments
            self.message_container.configure(style=f"{style['bootstyle']}.TFrame")
        else:
            # Light theme - use custom colors
            self.message_container.configure(style="TFrame")
            # We'll use the bootstyle colors which work well
        
        # Show the message
        self.message_container.pack(fill=X, pady=5)
        
        # Auto-hide after 5 seconds for non-error messages
        if status != "fail":
            self.root.after(5000, self.hide_message)
    
    def hide_message(self):
        """Hide the message bar"""
        self.message_container.pack_forget()
        
    def create_section(self, parent, title):
        """Create a styled section with title"""
        frame = ttk.LabelFrame(parent, text=title, padding=15)
        frame.pack(fill=BOTH, expand=TRUE, pady=(0, 15))
        return frame
    
    def select_folder(self):
        """Select working folder"""
        folder = filedialog.askdirectory()
        if folder:
            self.working_folder.set(folder)
            self.populate_file_list()
    
    def populate_file_list(self, show_toast=True):
        """Populate file listbox with supported files"""
        self.file_listbox.delete(0, tk.END)
        
        folder = self.working_folder.get()
        if folder and os.path.exists(folder):
            files_found = 0
            for file in os.listdir(folder):
                if file.lower().endswith(('.txt', '.pdf', '.docx', '.pptx')):
                    self.file_listbox.insert(tk.END, f"📄 {file}")
                    files_found += 1
            
            if files_found == 0:
                self.file_listbox.insert(tk.END, "No supported files found")
                if show_toast:
                    self.show_message("No supported files found in the selected folder", "info")
            else:
                if show_toast:
                    self.show_message(f"Found {files_found} supported file(s) in the selected folder", "info")
    
    def select_all_files(self):
        """Select all files in the listbox"""
        # First check if there are any files
        if self.file_listbox.size() == 0:
            self.show_message("No files to select", "warning")
            return
        
        # Check if the first item is "No supported files found"
        first_item = self.file_listbox.get(0)
        if first_item == "No supported files found":
            self.show_message("No files to select", "warning")
            return
        
        # Select all items
        self.file_listbox.selection_set(0, tk.END)
        
        # Show message
        file_count = self.file_listbox.size()
        self.show_message(f"Selected all {file_count} file(s)", "info")
    
    def clear_file_selection(self):
        """Clear all file selections"""
        self.file_listbox.selection_clear(0, tk.END)
        self.text_content = ""
        # Update source mode label
        self.update_mode_label()
        self.show_message("File selection cleared", "info")
    
    def load_files(self):
        """Load selected files"""
        selected_indices = self.file_listbox.curselection()
        if not selected_indices:
            self.show_message("Please select at least one file to load", "warning")
            return
        
        self.text_content = ""
        folder = self.working_folder.get()
        
        # Update source mode label
        self.update_mode_label(source="Files")
        
        for idx in selected_indices:
            filename = self.file_listbox.get(idx).replace("📄 ", "")
            filepath = os.path.join(folder, filename)
            
            try:
                if filename.lower().endswith('.txt'):
                    with open(filepath, 'r', encoding='utf-8') as f:
                        self.text_content += f.read() + "\n"
                
                elif filename.lower().endswith('.pdf'):
                    with open(filepath, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        for page in pdf_reader.pages:
                            page_text = page.extract_text()
                            # Fix common PDF extraction issues
                            # Replace soft hyphens and rejoin split words
                            page_text = page_text.replace('\u00AD', '')  # Remove soft hyphens
                            page_text = page_text.replace('-\n', '')  # Rejoin hyphenated words
                            page_text = page_text.replace('\n', ' ')  # Replace newlines with spaces
                            self.text_content += page_text + " "
                
                elif filename.lower().endswith('.docx'):
                    doc = Document(filepath)
                    for paragraph in doc.paragraphs:
                        self.text_content += paragraph.text + "\n"
                
                elif filename.lower().endswith('.pptx'):
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                self.text_content += shape.text + "\n"
                
            except Exception as e:
                self.show_message(f"Error reading {filename}: {str(e)}", "fail")
                self.show_toast(f"Error reading {filename}", "danger")
        
        # Show success message in the message bar
        total_words = len(self.text_content.split())
        self.show_message(f"Successfully loaded {len(selected_indices)} file(s) with approximately {total_words:,} words", "good")
        self.show_toast(f"Files loaded successfully!", "success")
    
    def use_pasted_text(self):
        """Use text from text input widget"""
        self.text_content = self.text_input.get('1.0', tk.END).strip()
        if self.text_content:
            # Update source mode label
            self.update_mode_label(source="Custom Text")
            word_count = len(self.text_content.split())
            self.show_message(f"Text loaded successfully with approximately {word_count:,} words", "good")
            self.show_toast("Text loaded successfully", "success")
        else:
            self.show_message("Please paste some text into the text area first", "warning")
            self.show_toast("Please paste some text first", "warning")
    
    def update_min_label(self, value):
        """Update minimum length label"""
        val = int(float(value))
        self.min_word_length.set(val)
        self.min_length_label.config(text=str(val))
        # Ensure max is not less than min
        if self.max_length_scale.get() < val:
            self.max_length_scale.set(val)
    
    def update_max_label(self, value):
        """Update maximum length label"""
        val = int(float(value))
        self.max_word_length.set(val)
        self.max_length_label.config(text=str(val))
        # Ensure min is not greater than max
        if self.min_length_scale.get() > val:
            self.min_length_scale.set(val)
    
    def update_min_from_meter(self):
        """Update min length from meter widget"""
        if self.min_length_meter:
            val = int(self.min_length_meter.amountusedvar.get())
            self.min_word_length.set(val)
            # Ensure max is not less than min
            if self.max_length_meter and self.max_length_meter.amountusedvar.get() < val:
                self.max_length_meter.amountusedvar.set(val)
    
    def update_max_from_meter(self):
        """Update max length from meter widget"""
        if self.max_length_meter:
            val = int(self.max_length_meter.amountusedvar.get())
            self.max_word_length.set(val)
            # Ensure min is not greater than max
            if self.min_length_meter and self.min_length_meter.amountusedvar.get() > val:
                self.min_length_meter.amountusedvar.set(val)
    
    def update_mode_label(self, source=None):
        """Update the mode label with source and mask information"""
        if not hasattr(self, 'source_label') or not hasattr(self, 'mask_label'):
            return
        
        # Determine source
        if source is None:
            if self.text_content:
                # Check if it's from files or custom text
                if hasattr(self, 'file_listbox') and self.file_listbox.curselection():
                    source = "Files"
                else:
                    source = "Custom Text"
            else:
                source = "None"
        
        # Source icons
        source_icons = {
            "None": "📄",
            "Files": "📁",
            "Custom Text": "✏️"
        }
        source_icon = source_icons.get(source, "📄")
        
        # Update source label
        if self.text_content and source != "None":
            word_count = len(self.text_content.split())
            source_text = f"{source} ({word_count:,} words)"
        else:
            source_text = source
        
        self.source_label.config(text=source_text)
        
        # Determine mask type with icons
        mask_icons = {
            "none": "⬜",
            "no_mask": "⬜",
            "image": "🖼️",
            "image_mask": "🖼️",
            "text": "🔤",
            "text_mask": "🔤"
        }
        
        # Get mask type from the radio button selection
        mask_type = self.mask_type.get() if hasattr(self, 'mask_type') else "no_mask"
        mask_icon = mask_icons.get(mask_type, "⬜")
        
        # Build mask description
        if mask_type in ["none", "no_mask"]:
            mask_desc = "No Mask"
        elif mask_type in ["image", "image_mask"]:
            mask_desc = "Image Mask"
        elif mask_type in ["text", "text_mask"]:
            mask_desc = "Text Mask"
        else:
            mask_desc = "No Mask"
        
        # Update mask label
        self.mask_label.config(text=mask_desc)
    
    def update_color_scheme_label(self):
        """Update the color scheme label based on current selection"""
        if not hasattr(self, 'color_scheme_label'):
            return
            
        mode = self.color_mode.get()
        
        if mode == "single":
            self.color_scheme_label.config(text="Single Color")
        elif mode == "preset":
            # Get the selected preset name
            preset_name = self.color_var.get()
            self.color_scheme_label.config(text=preset_name)
        elif mode == "custom":
            self.color_scheme_label.config(text="Custom Gradient")
    
    def update_status_bar_colors(self, bg_color, border_color, text_color, label_color):
        """Update status bar colors based on theme"""
        # Update main container
        self.status_bar.config(bg=bg_color)
        self.status_bar_inner.config(bg=bg_color)
        
        # Update borders
        self.status_bar_top_border.config(bg=border_color)
        
        # Find and update all child widgets
        def update_widget_colors(widget):
            try:
                # Skip certain widget types
                if isinstance(widget, (ttk.Label, ttk.Button, ttk.Frame)):
                    return
                    
                # Update background
                if hasattr(widget, 'config'):
                    widget.config(bg=bg_color)
                    
                    # Update text colors for labels
                    if isinstance(widget, tk.Label):
                        current_font = widget.cget('font')
                        if current_font and 'bold' in str(current_font):
                            widget.config(fg=text_color)
                        else:
                            widget.config(fg=label_color)
                    
                # Update divider
                if hasattr(widget, 'winfo_class') and widget.winfo_class() == 'Frame':
                    # Check if it's the divider by its width
                    try:
                        if widget.cget('width') == 1:
                            widget.config(bg=border_color)
                    except:
                        pass
                        
                # Recursively update children
                for child in widget.winfo_children():
                    update_widget_colors(child)
            except:
                pass
        
        update_widget_colors(self.status_bar)
        
        # Find and update bottom border
        for child in self.status_bar.winfo_children():
            try:
                if hasattr(child, 'cget') and child.cget('height') == 1 and child != self.status_bar_top_border:
                    child.config(bg=border_color)
                    break
            except:
                pass
    
    def update_forbidden_words(self, show_toast=True):
        """Update forbidden words set"""
        text = self.forbidden_text.get('1.0', tk.END).strip()
        # Only use the words explicitly listed in the text area, not STOPWORDS
        self.forbidden_words = set()
        if text:
            custom_forbidden = set(word.strip().lower() for word in text.split('\n') if word.strip())
            self.forbidden_words.update(custom_forbidden)
            self.print_debug(f"Custom forbidden words: {self.forbidden_words}")
        
        self.print_debug(f"Updated forbidden words from GUI text area: {len(self.forbidden_words)} words")
        if show_toast:
            self.show_toast(f"Updated forbidden words ({len(self.forbidden_words)} total)", "info")
    
    def on_color_select(self):
        """Handle color scheme selection"""
        color_name = self.color_var.get()
        self.selected_colormap = self.color_schemes[color_name]
        # Update combined preview if in preset mode
        if self.color_mode.get() == "preset":
            self.update_combined_color_preview()
        # Update color scheme label
        self.update_color_scheme_label()
        
    def on_color_tab_changed(self, event):
        """Handle color notebook tab change"""
        selected_tab = event.widget.index('current')
        if selected_tab == 0:  # Single Color
            self.color_mode.set("single")
        elif selected_tab == 1:  # Preset Gradients
            self.color_mode.set("preset")
        elif selected_tab == 2:  # Custom Gradient
            self.color_mode.set("custom")
        self.on_color_mode_change()
    
    def on_mask_tab_changed(self, event):
        """Handle mask notebook tab change"""
        selected_tab = event.widget.index('current')
        if selected_tab == 0:  # No Mask
            self.mask_type.set("no_mask")
        elif selected_tab == 1:  # Image Mask
            self.mask_type.set("image_mask")
        elif selected_tab == 2:  # Text Mask
            self.mask_type.set("text_mask")
        self.on_mask_type_change()
    
    def on_color_mode_change(self):
        """Handle color mode radio button change"""
        mode = self.color_mode.get()
        if mode == "single":
            self.color_notebook.select(0)
        elif mode == "preset":
            self.color_notebook.select(1)
        elif mode == "custom":
            self.color_notebook.select(2)
            self.update_custom_gradient_preview()
        
        # Update the combined preview
        self.update_combined_color_preview()
        
        # Update color scheme label
        self.update_color_scheme_label()
    
    def choose_custom_color(self, index):
        """Choose a color for custom gradient"""
        current_color = self.custom_gradient_colors[index]
        dialog = ColorChooserDialog(title=f"Choose Color {index+1}", 
                                   initialcolor=current_color)
        dialog.show()
        
        if dialog.result:
            color = dialog.result
            hex_color = color.hex
            self.custom_gradient_colors[index] = hex_color
            self.update_custom_gradient_preview()
    
    def add_gradient_color(self):
        """Add a new color to the gradient"""
        if len(self.custom_gradient_colors) < 10:  # Limit to 10 colors
            self.custom_gradient_colors.append("#FFFFFF")
            self.recreate_custom_gradient_ui()
    
    def remove_gradient_color(self):
        """Remove the last color from gradient"""
        if len(self.custom_gradient_colors) > 2:  # Minimum 2 colors
            self.custom_gradient_colors.pop()
            self.recreate_custom_gradient_ui()
    
    def recreate_custom_gradient_ui(self):
        """Recreate the custom gradient UI after adding/removing colors"""
        # Find the custom frame
        custom_tab = self.color_notebook.winfo_children()[2]
        custom_frame = custom_tab.winfo_children()[0]
        
        # Clear existing color frames
        for frame in self.custom_color_frames:
            frame.destroy()
        self.custom_color_frames.clear()
        self.custom_color_previews.clear()
        
        # Recreate color rows
        for i in range(len(self.custom_gradient_colors)):
            color_row = ttk.Frame(custom_frame)
            color_row.pack(fill=X, pady=5, before=custom_frame.winfo_children()[1])
            
            ttk.Label(color_row, text=f"Color {i+1}:", width=10).pack(side=LEFT)
            
            preview = ttk.Frame(color_row, width=30, height=30)
            preview.pack(side=LEFT, padx=(5, 10))
            self.custom_color_previews.append(preview)
            
            btn = ttk.Button(color_row, text="Choose", 
                           command=lambda idx=i: self.choose_custom_color(idx),
                           bootstyle="secondary-outline")
            btn.pack(side=LEFT)
            
            self.custom_color_frames.append(color_row)
        
        self.update_custom_gradient_preview()
    
    def update_custom_gradient_preview(self):
        """Update the custom gradient preview"""
        # Update individual color previews
        for i, (preview, color) in enumerate(zip(self.custom_color_previews, self.custom_gradient_colors)):
            style = ttk.Style()
            style_name = f"CustomColor{i}.TFrame"
            style.configure(style_name, background=color)
            preview.configure(style=style_name)
        
        # Update combined preview if in custom mode
        if self.color_mode.get() == "custom":
            self.update_combined_color_preview()
            
    def choose_single_color(self):
        """Open color picker for single color"""
        dialog = ColorChooserDialog(title="Choose Single Color", 
                                   initialcolor=self.single_color.get())
        dialog.show()
        
        if dialog.result:
            color = dialog.result
            hex_color = color.hex
            self.single_color.set(hex_color)
            
            # Update preview
            style = ttk.Style()
            style_name = "SingleColorPreview.TFrame"
            style.configure(style_name, background=hex_color)
            self.single_color_preview.configure(style=style_name)
            
            # Update combined preview if in single color mode
            if self.color_mode.get() == "single":
                self.update_combined_color_preview()
    
    
    def update_combined_color_preview(self):
        """Update the combined color preview based on selected mode"""
        if not hasattr(self, 'combined_color_preview'):
            return
            
        # Clear existing preview
        for widget in self.combined_color_preview.winfo_children():
            widget.destroy()
            
        mode = self.color_mode.get()
        
        if mode == "single":
            # Show single color
            preview_canvas = tk.Canvas(self.combined_color_preview, 
                                     height=50, 
                                     highlightthickness=0)
            preview_canvas.pack(fill=X)
            preview_canvas.configure(bg=self.single_color.get())
            
        elif mode == "preset":
            # Show selected preset gradient
            color_name = self.color_var.get()
            cmap_name = self.color_schemes.get(color_name)
            
            if cmap_name:
                preview_canvas = tk.Canvas(self.combined_color_preview, 
                                         height=50, 
                                         highlightthickness=0)
                preview_canvas.pack(fill=X)
                
                # Draw gradient when canvas is configured
                def draw_gradient(event=None):
                    width = preview_canvas.winfo_width()
                    if width > 1:
                        try:
                            # Get the actual colormap
                            cmap = matplotlib.colormaps[cmap_name]
                            for i in range(width):
                                color = cmap(i / width)
                                rgb = tuple(int(c * 255) for c in color[:3])
                                hex_color = '#%02x%02x%02x' % rgb
                                preview_canvas.create_line(i, 0, i, 50, fill=hex_color)
                        except:
                            pass
                
                preview_canvas.bind('<Configure>', draw_gradient)
                
        elif mode == "custom":
            # Show custom gradient
            if len(self.custom_gradient_colors) >= 2:
                preview_canvas = tk.Canvas(self.combined_color_preview, 
                                         height=50, 
                                         highlightthickness=0)
                preview_canvas.pack(fill=X)
                
                # Draw custom gradient when canvas is configured
                def draw_custom_gradient(event=None):
                    width = preview_canvas.winfo_width()
                    if width > 1:
                        try:
                            # Create custom colormap
                            custom_cmap = LinearSegmentedColormap.from_list(
                                'custom_gradient', 
                                self.custom_gradient_colors
                            )
                            
                            for i in range(width):
                                color = custom_cmap(i / width)
                                rgb = tuple(int(c * 255) for c in color[:3])
                                hex_color = '#%02x%02x%02x' % rgb
                                preview_canvas.create_line(i, 0, i, 50, fill=hex_color)
                        except:
                            pass
                
                preview_canvas.bind('<Configure>', draw_custom_gradient)
    
    def select_mask(self):
        """Select mask image file"""
        file_path = filedialog.askopenfilename(
            filetypes=[("Image files", "*.png *.jpg *.jpeg *.bmp *.gif")]
        )
        if file_path:
            try:
                self.image_mask_image = np.array(Image.open(file_path))
                self.mask_image = self.image_mask_image  # For backward compatibility
                self.mask_path.set(os.path.basename(file_path))
                
                # Update the image mask label
                self.image_mask_label.config(text=os.path.basename(file_path))
                
                # Update mask preview
                img = Image.open(file_path)
                img.thumbnail((200, 200), Image.Resampling.LANCZOS)
                photo = ImageTk.PhotoImage(img)
                if hasattr(self, 'image_mask_preview_label'):
                    self.image_mask_preview_label.config(image=photo, text="")
                    self.image_mask_preview_label.image = photo  # Keep a reference
                
                # Enable outline options when mask is selected
                self.update_outline_state(True)
                
                # Update mode label
                self.update_mode_label()
            except Exception as e:
                self.show_toast(f"Error loading mask: {str(e)}", "danger")
    
    def clear_mask(self):
        """Clear selected mask"""
        # Clear the appropriate mask based on current tab
        current_tab = self.mask_notebook.index(self.mask_notebook.select())
        
        if current_tab == 1:  # Image mask tab
            self.image_mask_image = None
            self.image_mask_label.config(text="No image selected")
        elif current_tab == 2:  # Text mask tab
            self.text_mask_image = None
            self.text_mask_input.set("")
            
        self.mask_image = None  # For backward compatibility
        self.mask_path.set("No mask selected")
        
        # Clear appropriate preview label
        if current_tab == 1 and hasattr(self, 'image_mask_preview_label'):
            self.image_mask_preview_label.config(image="", text="No mask selected")
        elif current_tab == 2 and hasattr(self, 'text_mask_preview_label'):
            self.text_mask_preview_label.config(image="", text="No mask selected")
        
        # Disable outline options when mask is cleared
        self.update_outline_state(False)
        
        # Update mode label
        self.update_mode_label()
    
    def create_text_mask(self, text, width=None, height=None, font_size=None):
        """Create a mask image from text"""
        if not text:
            return None
        
        # Use canvas dimensions if not specified
        if width is None:
            width = self.canvas_width.get()
        if height is None:
            height = self.canvas_height.get()
        if font_size is None:
            font_size = self.text_mask_font_size.get()
        
        # Create white image
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Get selected font
        selected_font = self.text_mask_font.get()
        font_name = self.available_fonts.get(selected_font, "Arial")
        
        # Build font style
        font_style = []
        if self.text_mask_bold.get():
            font_style.append("Bold")
        if self.text_mask_italic.get():
            font_style.append("Italic")
        
        # Try to load the font with style
        font = None
        font_attempts = []
        
        # First try with full style
        if font_style:
            font_attempts.append(f"{font_name} {' '.join(font_style)}")
        
        # Then try just the font name
        font_attempts.append(font_name)
        
        # Then try with .ttf extension
        font_attempts.append(f"{font_name.lower().replace(' ', '')}.ttf")
        
        # Try each font attempt
        for attempt in font_attempts:
            try:
                font = ImageFont.truetype(attempt, font_size)
                break
            except:
                continue
        
        # Fallback fonts
        if font is None:
            fallback_fonts = ["arial.ttf", "Arial", "helvetica", "verdana"]
            for fallback in fallback_fonts:
                try:
                    font = ImageFont.truetype(fallback, font_size)
                    break
                except:
                    continue
        
        # Final fallback to default
        if font is None:
            font = ImageFont.load_default()
        
        # Handle multi-line text
        words_per_line = self.text_mask_words_per_line.get()
        if words_per_line > 1:
            # Split text into words and group them
            words = text.split()
            lines = []
            for i in range(0, len(words), words_per_line):
                lines.append(' '.join(words[i:i+words_per_line]))
            text_to_draw = '\n'.join(lines)
        else:
            text_to_draw = text
        
        # Get text boundaries for multi-line text
        bbox = draw.textbbox((0, 0), text_to_draw, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        # Center the text
        x = (width - text_width) // 2
        y = (height - text_height) // 2
        
        # Draw text in black with stroke thickness and letter spacing
        # Map thickness value (0-5) to stroke width (0-15 pixels for more visible effect)
        thickness_val = self.letter_thickness.get()
        stroke_width = int(thickness_val * 3) if thickness_val > 0 else 0
        spacing = int(self.letter_spacing.get() * 10)  # Convert 0-5 to 0-50 pixels
        
        # If no spacing needed, use normal drawing
        if spacing == 0:
            try:
                # Try using stroke_width (requires PIL 6.2.0+)
                draw.text((x, y), text_to_draw, fill='black', font=font, align='center', 
                          stroke_width=stroke_width, stroke_fill='black')
            except TypeError:
                # Fallback for older PIL versions - just draw without stroke
                draw.text((x, y), text_to_draw, fill='black', font=font, align='center')
        else:
            # Draw text with letter spacing - handle each line separately
            lines = text_to_draw.split('\n')
            current_y = y
            
            for line in lines:
                # Calculate line width with spacing
                line_width = 0
                for char in line:
                    bbox = draw.textbbox((0, 0), char, font=font)
                    char_width = bbox[2] - bbox[0]
                    line_width += char_width + spacing
                line_width -= spacing  # Remove last spacing
                
                # Start position for this line (centered)
                current_x = x + (text_width - line_width) // 2
                
                # Draw each character with spacing
                for char in line:
                    try:
                        # Try with stroke
                        draw.text((current_x, current_y), char, fill='black', font=font,
                                  stroke_width=stroke_width, stroke_fill='black')
                    except TypeError:
                        # Fallback without stroke
                        draw.text((current_x, current_y), char, fill='black', font=font)
                    
                    # Move to next character position
                    bbox = draw.textbbox((0, 0), char, font=font)
                    char_width = bbox[2] - bbox[0]
                    current_x += char_width + spacing
                
                # Move to next line
                bbox = draw.textbbox((0, 0), 'Ay', font=font)  # Sample for line height
                line_height = bbox[3] - bbox[1]
                current_y += line_height
        
        # Convert to numpy array
        return np.array(img)
    
    def update_text_mask(self):
        """Update the text mask when text or settings change"""
        if self.text_mask_input.get():
            # Generate text mask
            self.text_mask_image = self.create_text_mask(self.text_mask_input.get())
            self.mask_image = self.text_mask_image  # For backward compatibility
            self.mask_path.set(f"Text: {self.text_mask_input.get()}")
            
            # Update preview
            self.update_mask_preview()
            
            # Enable outline options
            self.update_outline_state(True)
            
            # Update mode label
            self.update_mode_label()
    
    def update_mask_preview(self):
        """Update the mask preview display"""
        # Determine which mask to preview based on context
        mask_to_preview = None
        preview_label = None
        
        # Check which tab is active
        current_tab = self.mask_notebook.index(self.mask_notebook.select())
        
        if current_tab == 1 and self.image_mask_image is not None:  # Image mask tab
            mask_to_preview = self.image_mask_image
            if hasattr(self, 'image_mask_preview_label'):
                preview_label = self.image_mask_preview_label
        elif current_tab == 2 and self.text_mask_image is not None:  # Text mask tab
            mask_to_preview = self.text_mask_image
            if hasattr(self, 'text_mask_preview_label'):
                preview_label = self.text_mask_preview_label
                
        # Clear all preview labels first
        if hasattr(self, 'text_mask_preview_label'):
            self.text_mask_preview_label.config(image="", text="")
            if hasattr(self.text_mask_preview_label, 'image'):
                self.text_mask_preview_label.image = None
        if hasattr(self, 'image_mask_preview_label'):
            self.image_mask_preview_label.config(image="", text="")
            if hasattr(self.image_mask_preview_label, 'image'):
                self.image_mask_preview_label.image = None
                
        if mask_to_preview is not None and preview_label is not None:
            # Convert numpy array to PIL Image for preview
            if len(mask_to_preview.shape) == 3:
                preview_img = Image.fromarray(mask_to_preview.astype('uint8'), 'RGB')
            else:
                preview_img = Image.fromarray(mask_to_preview.astype('uint8'), 'L')
            
            # Calculate preview size based on canvas dimensions
            canvas_width = self.canvas_width.get()
            canvas_height = self.canvas_height.get()
            preview_width = int(canvas_width * 0.25)  # 25% of canvas width
            preview_height = int(canvas_height * 0.25)  # 25% of canvas height
            
            # Resize for preview maintaining aspect ratio
            preview_img.thumbnail((preview_width, preview_height), Image.Resampling.LANCZOS)
            
            # Add border to text mask preview
            if current_tab == 2:  # Text mask tab
                # Create a new image with border
                border_width = 2
                border_color = '#6B7280'  # Gray border
                
                # Get the size after thumbnail
                img_width, img_height = preview_img.size
                
                # Create new image with border
                bordered_img = Image.new('RGB', 
                                       (img_width + 2*border_width, img_height + 2*border_width), 
                                       border_color)
                
                # Paste the preview image in the center
                bordered_img.paste(preview_img, (border_width, border_width))
                preview_img = bordered_img
            
            photo = ImageTk.PhotoImage(preview_img)
            
            # Update the preview label
            preview_label.config(image=photo, text="")
            preview_label.image = photo
    
    def update_outline_width(self, value, label=None):
        """Update outline width label"""
        val = int(float(value))
        self.outline_width.set(val)
        if label:
            label.config(text=f"{val} pixels")
        elif hasattr(self, 'outline_width_label'):
            self.outline_width_label.config(text=f"{val} pixels")
    
    def choose_outline_color(self, preview_frame=None):
        """Open color chooser for outline color"""
        dialog = ColorChooserDialog()
        dialog.show()
        color = dialog.result
        if color:
            hex_color = color.hex
            self.outline_color.set(hex_color)
            # Update preview - ttk frames don't support background, use style instead
            style = ttk.Style()
            style_name = f"OutlinePreview.TFrame"
            style.configure(style_name, background=hex_color)
            if preview_frame:
                preview_frame.configure(style=style_name)
            elif hasattr(self, 'outline_color_preview'):
                self.outline_color_preview.configure(style=style_name)
    
    def choose_bg_color(self):
        """Open color chooser for background color"""
        dialog = ColorChooserDialog()
        dialog.show()
        color = dialog.result
        if color:
            hex_color = color.hex
            self.bg_color.set(hex_color)
            self.update_bg_preview()
    
    def update_bg_preview(self):
        """Update the background color preview"""
        if hasattr(self, 'bg_color_preview'):
            # Update preview - ttk frames don't support background, use style instead
            style = ttk.Style()
            style_name = f"BgPreview.TFrame"
            style.configure(style_name, background=self.bg_color.get())
            self.bg_color_preview.configure(style=style_name)
    
    def on_color_mode_change_canvas(self):
        """Handle canvas color mode change between RGB and RGBA"""
        if hasattr(self, 'bg_color_btn') and hasattr(self, 'rgba_mode'):
            if self.rgba_mode.get():
                # RGBA mode - disable background color selection
                self.bg_color_btn.configure(state='disabled')
                if hasattr(self, 'bg_color_preview'):
                    self.bg_color_preview.configure(bootstyle="secondary")
            else:
                # RGB mode - enable background color selection
                self.bg_color_btn.configure(state='normal')
                self.update_bg_preview()
    
    def update_outline_color_preview(self):
        """Update the outline color preview"""
        if hasattr(self, 'outline_color_preview') and hasattr(self, 'outline_color'):
            style = ttk.Style()
            style_name = f"OutlinePreview{id(self)}.TFrame"
            style.configure(style_name, background=self.outline_color)
            self.outline_color_preview.configure(style=style_name)
    
    def clear_canvas(self, clear_wordcloud=True, show_placeholder=True):
        """Clear the canvas completely"""
        self.print_debug("Clearing canvas...")
        
        # Store current axes for complete removal
        axes_to_remove = self.figure.axes[:]
        
        # Remove all axes completely
        for ax in axes_to_remove:
            ax.clear()
            self.figure.delaxes(ax)
        
        # Clear the figure completely
        self.figure.clear()
        
        # Force garbage collection of matplotlib objects
        import gc
        gc.collect()
        
        # Create fresh subplot
        ax = self.figure.add_subplot(111)
        ax.set_facecolor('white')
        if show_placeholder:
            ax.text(0.5, 0.5, '', 
                    horizontalalignment='center', verticalalignment='center',
                    transform=ax.transAxes, fontsize=14, color='gray')
        ax.axis('off')
        
        # Reset figure properties based on theme
        if hasattr(self, 'current_theme') and self.current_theme.get() in ["darkly", "superhero", "solar", "cyborg", "vapor"]:
            self.figure.patch.set_facecolor('#2b2b2b')
        else:
            self.figure.patch.set_facecolor('white')
        
        # Clear the canvas widget
        self.canvas.draw_idle()
        self.canvas.flush_events()
        
        # Process pending GUI events to ensure complete update
        self.root.update_idletasks()
        
        # Disable save button since there's nothing to save
        if hasattr(self, 'save_btn'):
            self.save_btn.config(state=DISABLED)
        
        # Clear wordcloud object only if requested
        if clear_wordcloud:
            if hasattr(self, 'wordcloud'):
                self.wordcloud = None
            if hasattr(self, 'current_wordcloud'):
                self.current_wordcloud = None
            
        self.print_debug("Canvas cleared successfully")
    
    def update_preview_size(self, *args):
        """Update preview canvas size when dimensions change"""
        try:
            # Calculate new display size
            display_width, display_height = self.calculate_preview_size()
            
            # Update figure size for display
            self.figure.set_size_inches(display_width/100, display_height/100)
            
            # Update the canvas widget size to match
            self.canvas_widget.config(width=display_width, height=display_height)
            
            # Update scale indicator and slider
            actual_width = self.canvas_width.get()
            actual_height = self.canvas_height.get()
            
            # Calculate actual scale percentage
            if actual_width > 0:
                actual_scale = int((display_width / actual_width) * 100)
                
                # Update slider and label if they exist
                if hasattr(self, 'preview_slider'):
                    self.preview_scale.set(actual_scale)
                    self.preview_slider.set(actual_scale)
                if hasattr(self, 'preview_size_label'):
                    self.preview_size_label.config(text=f"{actual_scale}%")
                
                # Update scale indicator
                if display_width < actual_width or display_height < actual_height:
                    self.scale_indicator.config(text=f"Preview at {actual_scale}% (limited by screen size)")
                else:
                    self.scale_indicator.config(text="")
            
            # Don't clear canvas when preview size changes
            # self.clear_canvas()
            
            # Force canvas to redraw with new size
            self.canvas.draw()
        except Exception as e:
            if hasattr(self, 'print_debug'):
                self.print_debug(f"Error updating preview size: {str(e)}")
    
    def adjust_preview_size(self, delta):
        """Adjust preview size by delta percent"""
        current = self.preview_scale.get()
        new_value = max(25, min(200, current + delta))
        self.preview_scale.set(new_value)
        self.preview_slider.set(new_value)  # Sync the slider
        self.update_preview_size_from_slider(new_value)
    
    def update_preview_size_from_slider(self, value):
        """Update preview size from slider value"""
        val = int(float(value))
        self.preview_scale.set(val)
        self.preview_size_label.config(text=f"{val}%")
        
        # Update the preview canvas
        self.update_preview_display()
    
    def set_preview_size(self, percent):
        """Set preview size to specific percentage"""
        self.preview_scale.set(percent)
        self.preview_slider.set(percent)  # Sync the slider
        self.update_preview_size_from_slider(percent)
    
    def update_preview_display(self):
        """Update the preview display based on current scale"""
        try:
            scale_factor = self.preview_scale.get() / 100.0
            
            # Get base dimensions
            base_width = self.canvas_width.get()
            base_height = self.canvas_height.get()
            
            # Calculate scaled dimensions
            scaled_width = int(base_width * scale_factor)
            scaled_height = int(base_height * scale_factor)
            
            # Update figure size
            self.figure.set_size_inches(scaled_width/100, scaled_height/100)
            
            # Update canvas widget size
            self.canvas_widget.config(width=scaled_width, height=scaled_height)
            
            # Update scale indicator
            display_width, display_height = self.calculate_preview_size()
            actual_width = self.canvas_width.get()
            actual_height = self.canvas_height.get()
            
            # Check if constrained by screen size
            user_scale = self.preview_scale.get() / 100.0
            if display_width < actual_width * user_scale or display_height < actual_height * user_scale:
                # Constrained by screen limits
                actual_percent = int((display_width / actual_width) * 100)
                self.scale_indicator.config(text=f"Preview at {actual_percent}% (limited by screen size)")
            elif scale_factor < 1.0:
                self.scale_indicator.config(text=f"Preview at {self.preview_scale.get()}% of actual size")
            elif scale_factor > 1.0:
                self.scale_indicator.config(text=f"Preview enlarged to {self.preview_scale.get()}%")
            else:
                self.scale_indicator.config(text="Preview at actual size")
            
            # Redraw canvas
            self.canvas.draw()
            
            # If there's a wordcloud, regenerate it at the new size
            if hasattr(self, 'current_wordcloud') and self.current_wordcloud:
                self.display_wordcloud(self.current_wordcloud)
                
        except Exception as e:
            self.print_debug(f"Error updating preview display: {str(e)}")
    
    def update_outline_state(self, has_mask=None):
        """Enable/disable outline options based on mask selection"""
        if has_mask is None:
            # Check if any mask is selected based on radio button
            mask_type = self.mask_type.get()
            if mask_type == "image_mask":
                has_mask = self.image_mask_image is not None
            elif mask_type == "text_mask":
                has_mask = self.text_mask_image is not None
            else:
                has_mask = False
        
        state = NORMAL if has_mask else DISABLED
        
        for widget in self.outline_widgets:
            try:
                widget.configure(state=state)
            except:
                pass  # Some widgets might not support state
    
    def update_horizontal_label(self, value):
        """Update prefer horizontal label"""
        val = float(value)
        self.prefer_horizontal.set(val)
        self.horizontal_label.config(text=f"{int(val * 100)}%")
    
    def update_horizontal_gauge(self, value):
        """Update horizontal gauge and prefer_horizontal value"""
        val = float(value)
        self.horizontal_gauge.configure(value=val)
        self.prefer_horizontal.set(val / 100.0)  # Convert percentage to 0-1 range
    
    def reset_orientation(self):
        """Reset word orientation to default 90%"""
        self.horizontal_scale.set(90)
        self.horizontal_gauge.configure(value=90)
        self.prefer_horizontal.set(0.9)
        self.show_toast("Word orientation reset to 90% horizontal", "info")
        # Don't clear canvas on orientation reset
        # self.clear_canvas()
    
    def update_thickness_label(self, value):
        """Update letter thickness label"""
        val = float(value)
        self.letter_thickness.set(val)
        if val < 0.5:
            text = "Very Thin"
        elif val < 0.8:
            text = "Thin"
        elif val < 1.2:
            text = "Normal"
        elif val < 2.0:
            text = "Thick"
        else:
            text = "Very Thick"
        self.thickness_label.config(text=text)
        # Update text mask if it's currently selected
        if hasattr(self, 'mask_type') and self.mask_type.get() == "text_mask":
            self.update_text_mask()
    
    def update_max_words(self, value):
        """Update max words label"""
        val = int(float(value))
        self.max_words.set(val)
        if hasattr(self, 'max_words_label'):
            self.max_words_label.config(text=str(val))
        # Don't clear canvas on max words change
        # self.clear_canvas()
    
    def update_max_words_from_meter(self):
        """Update max words from meter widget"""
        if self.max_words_meter:
            val = int(self.max_words_meter.amountusedvar.get())
            self.max_words.set(val)
            # Don't clear canvas on max words change
            # self.clear_canvas()
    
    def update_scale(self, value):
        """Update scale label"""
        val = int(float(value))
        self.scale.set(val)
        if hasattr(self, 'scale_label'):
            self.scale_label.config(text=str(val))
        # Don't clear canvas on scale change
        # self.clear_canvas()
    
    def update_scale_from_meter(self):
        """Update scale from meter widget"""
        if self.scale_meter:
            val = int(self.scale_meter.amountusedvar.get())
            self.scale.set(val)
            # Don't clear canvas on scale change
            # self.clear_canvas()
    
    def update_thickness_from_meter(self):
        """Update letter thickness from meter widget"""
        if hasattr(self, 'thickness_meter') and self.thickness_meter:
            val = self.thickness_meter.amountusedvar.get()
            self.letter_thickness.set(val)
            # Update text mask if it's currently selected
            if hasattr(self, 'mask_type') and self.mask_type.get() == "text_mask":
                self.update_text_mask()
    
    def update_spacing_from_meter(self):
        """Update letter spacing from meter widget"""
        if hasattr(self, 'spacing_meter') and self.spacing_meter:
            val = self.spacing_meter.amountusedvar.get()
            self.letter_spacing.set(val)
            # Update text mask if it's currently selected
            if hasattr(self, 'mask_type') and self.mask_type.get() == "text_mask":
                self.update_text_mask()
    
    def update_words_per_line_from_meter(self):
        """Update words per line from meter widget"""
        if self.words_per_line_meter:
            val = int(self.words_per_line_meter.amountusedvar.get())
            self.text_mask_words_per_line.set(val)
            if self.text_mask_input.get():
                self.update_text_mask()
    
    def update_outline_width_from_meter(self):
        """Update outline width from meter widget"""
        if self.outline_width_meter:
            try:
                raw_val = self.outline_width_meter.amountusedvar.get()
                val = int(raw_val)
                # Force to 0 if less than 1
                if raw_val < 1:
                    val = 0
                self.outline_width.set(val)
                if self.text_mask_preview_label and hasattr(self.text_mask_preview_label, 'original_image'):
                    self.update_text_mask()
            except Exception as e:
                # Silently ignore errors during rapid updates
                pass
    
    def update_font_size_from_meter(self):
        """Update font size from meter widget"""
        if self.font_size_meter:
            val = int(self.font_size_meter.amountusedvar.get())
            self.text_mask_font_size.set(val)
            if self.text_mask_input.get():
                self.update_text_mask()
    
    def update_width_from_meter(self):
        """Update width from meter widget"""
        if self.width_meter:
            val = int(self.width_meter.amountusedvar.get())
            self.canvas_width.set(val)
            
            if self.lock_aspect_ratio.get() and self.aspect_ratio > 0:
                # Update height to maintain aspect ratio
                new_height = int(val / self.aspect_ratio)
                new_height = max(300, min(4000, new_height))
                self.canvas_height.set(new_height)
                if self.height_meter:
                    self.height_meter.amountusedvar.set(new_height)
            
            self.clear_canvas()
            # Update text mask preview if active
            if hasattr(self, 'mask_type') and self.mask_type.get() == "text_mask":
                self.update_text_mask()
    
    def update_height_from_meter(self):
        """Update height from meter widget"""
        if self.height_meter:
            val = int(self.height_meter.amountusedvar.get())
            self.canvas_height.set(val)
            
            if self.lock_aspect_ratio.get() and self.aspect_ratio > 0:
                # Update width to maintain aspect ratio
                new_width = int(val * self.aspect_ratio)
                new_width = max(400, min(4000, new_width))
                self.canvas_width.set(new_width)
                if self.width_meter:
                    self.width_meter.amountusedvar.set(new_width)
            
            self.clear_canvas()
            # Update text mask preview if active
            if hasattr(self, 'mask_type') and self.mask_type.get() == "text_mask":
                self.update_text_mask()
    
    def update_mode(self, show_toast=True):
        """Update mode between RGB and RGBA"""
        if self.rgba_mode.get():
            # RGBA mode - disable background color
            self.bg_label.configure(state=DISABLED)
            self.bg_color_btn.configure(state=DISABLED)
            
            # Show transparency checkbox
            if hasattr(self, 'transparency_check'):
                self.transparency_check.configure(state=NORMAL)
            
            # Disable outline options in RGBA mode
            if hasattr(self, 'outline_width_scale'):
                self.outline_width_scale.configure(state=DISABLED)
            if hasattr(self, 'outline_color_btn'):
                self.outline_color_btn.configure(state=DISABLED)
            
            if show_toast:
                self.show_toast("RGBA mode enabled - background will be transparent", "info")
                if self.outline_width.get() > 0:
                    self.show_toast("Note: Outlines disabled in RGBA mode", "warning")
        else:
            # RGB mode - enable background color
            self.bg_label.configure(state=NORMAL)
            self.bg_color_btn.configure(state=NORMAL)
            
            # Hide transparency checkbox
            if hasattr(self, 'transparency_check'):
                self.transparency_check.configure(state=DISABLED)
            
            # Re-enable outline options if mask is selected
            if hasattr(self, 'mask_image') and self.mask_image is not None:
                if hasattr(self, 'outline_width_scale'):
                    self.outline_width_scale.configure(state=NORMAL)
                if hasattr(self, 'outline_color_btn'):
                    self.outline_color_btn.configure(state=NORMAL)
            
            if show_toast:
                self.show_toast("RGB mode enabled - solid background", "info")
    
    def filter_words(self, text):
        """Filter words based on length and forbidden words"""
        # Clean up text first
        # Remove extra spaces and normalize whitespace
        text = ' '.join(text.split())
        
        # Extract words - include apostrophes for contractions
        words = re.findall(r"\b[\w']+\b", text.lower())
        
        # Additional cleanup - remove standalone punctuation and numbers
        words = [w for w in words if not w.isdigit() and len(w) > 0]
        
        # Update forbidden words (don't show toast during generation)
        self.update_forbidden_words(show_toast=False)
        
        # Filter words
        filtered_words = []
        min_len = self.min_word_length.get()
        max_len = self.max_word_length.get()
        
        self.print_debug(f"Filtering words: min_length={min_len}, max_length={max_len}, total_words={len(words)}")
        
        # Count words by length for debugging
        length_counts = {}
        filtered_by_length = 0
        filtered_by_forbidden = 0
        
        # Debug: show first 10 words being processed
        debug_limit = 10
        words_shown = 0
        
        for word in words:
            word_len = len(word)
            length_counts[word_len] = length_counts.get(word_len, 0) + 1
            
            # Detailed debug for first few words
            if words_shown < debug_limit:
                if min_len <= word_len <= max_len:
                    if word not in self.forbidden_words:
                        self.print_debug(f"  ✓ '{word}' (len={word_len}) - KEPT")
                        filtered_words.append(word)
                    else:
                        self.print_debug(f"  ✗ '{word}' (len={word_len}) - FORBIDDEN")
                        filtered_by_forbidden += 1
                else:
                    self.print_debug(f"  ✗ '{word}' (len={word_len}) - TOO SHORT/LONG")
                    filtered_by_length += 1
                words_shown += 1
            else:
                # Just count for remaining words
                if min_len <= word_len <= max_len:
                    if word not in self.forbidden_words:
                        filtered_words.append(word)
                    else:
                        filtered_by_forbidden += 1
                else:
                    filtered_by_length += 1
        
        # Log length distribution for words under min_length
        short_words = {k: v for k, v in length_counts.items() if k < min_len}
        if short_words:
            self.print_debug(f"Words shorter than min_length ({min_len}): {short_words}")
        
        self.print_debug(f"After filtering: {len(filtered_words)} words remain")
        self.print_debug(f"Filtered out: {filtered_by_length} by length, {filtered_by_forbidden} by forbidden list")
        
        return ' '.join(filtered_words)
    
    def validate_configuration(self):
        """Validate configuration and return list of warnings/errors"""
        issues = []
        
        # Check if min length > max length
        if self.min_word_length.get() > self.max_word_length.get():
            issues.append(("error", "Minimum word length cannot be greater than maximum word length"))
        
        # Check canvas size
        width = self.canvas_width.get()
        height = self.canvas_height.get()
        if width < 100 or height < 100:
            issues.append(("error", "Canvas size too small. Minimum size is 100x100"))
        if width > 4000 or height > 4000:
            issues.append(("warning", "Large canvas size may cause slow generation"))
        
        # Check if using text mask with no text
        if self.mask_type.get() == "text_mask" and (not hasattr(self, 'text_mask_image') or self.text_mask_image is None):
            issues.append(("error", "Text mask selected but no text provided"))
        
        # Check if using image mask with no image
        if self.mask_type.get() == "image_mask" and (not hasattr(self, 'image_mask_image') or self.image_mask_image is None):
            issues.append(("error", "Image mask selected but no image loaded"))
        
        # Check if using RGBA mode with outlines
        has_mask = False
        if self.mask_type.get() == "image_mask" and hasattr(self, 'image_mask_image') and self.image_mask_image is not None:
            has_mask = True
        elif self.mask_type.get() == "text_mask" and hasattr(self, 'text_mask_image') and self.text_mask_image is not None:
            has_mask = True
            
        if self.rgba_mode.get() and has_mask and self.outline_width.get() > 0:
            issues.append(("error", "Outlines are not supported in RGBA (transparent) mode. Please disable outlines or switch to RGB mode"))
        
        # Check max words
        if self.max_words.get() < 5:
            issues.append(("warning", "Very few words selected. Word cloud may look sparse"))
        
        # Check scale value
        if self.scale.get() > 5:
            issues.append(("warning", "High scale value may cause very slow generation"))
        
        # Check if all words might be filtered
        text_preview = self.text_content[:1000] if self.text_content else ""
        if text_preview:
            words = text_preview.split()
            avg_word_length = sum(len(w) for w in words) / len(words) if words else 0
            if avg_word_length < self.min_word_length.get():
                issues.append(("warning", f"Average word length ({avg_word_length:.1f}) is less than minimum filter ({self.min_word_length.get()}). Most words may be filtered out"))
        
        return issues
    
    def generate_wordcloud(self):
        """Generate word cloud in a separate thread"""
        if not self.text_content:
            self.show_message("No text content available. Please load files or paste text first.", "warning")
            self.show_toast("Please load text from files or paste text first", "warning")
            return
        
        # Validate configuration
        issues = self.validate_configuration()
        if issues:
            # Show errors first
            errors = [msg for level, msg in issues if level == "error"]
            warnings = [msg for level, msg in issues if level == "warning"]
            
            if errors:
                error_msg = "Cannot generate word cloud:\n\n" + "\n".join(f"• {msg}" for msg in errors)
                self.show_message(error_msg, "error")
                return
            
            if warnings:
                warning_msg = "Warnings:\n" + "\n".join(f"• {msg}" for msg in warnings)
                self.show_toast(warning_msg, "warning")
        
        # Clear the canvas before generating new word cloud
        self.print_debug("Clearing canvas before generation")
        # Force a complete clear before generation
        self.clear_canvas()
        # Additional forced update
        self.canvas.draw()
        self.root.update_idletasks()
        
        # Show progress and disable button
        self.generate_btn.config(state=DISABLED)
        self.progress.pack(fill=X, pady=(0, 10))
        self.progress.start(10)
        
        # Run generation in separate thread
        thread = threading.Thread(target=self._generate_wordcloud_thread)
        thread.start()
    
    def _generate_wordcloud_thread(self):
        """Generate word cloud (thread function)"""
        try:
            # Filter words
            filtered_text = self.filter_words(self.text_content)
            
            if not filtered_text:
                self.root.after(0, lambda: self.show_toast("No words found after filtering", "warning"))
                return
            
            # Create word cloud
            wc_params = {
                'width': self.canvas_width.get(),
                'height': self.canvas_height.get(),
                'max_words': int(self.max_words.get()),
                'scale': self.scale.get(),
                'relative_scaling': 0.5,
                'min_font_size': 4,  # Reduced from 10 to allow smaller words in masks
                'prefer_horizontal': self.prefer_horizontal.get(),
                'margin': int(5 * self.letter_thickness.get())  # Margin affects letter thickness
            }
            
            # Set color mode
            if self.color_mode.get() == "single":
                # Use single color function
                color_value = self.single_color.get()
                wc_params['color_func'] = lambda *args, **kwargs: color_value
            elif self.color_mode.get() == "custom":
                # Use custom gradient
                custom_cmap = LinearSegmentedColormap.from_list('custom', self.custom_gradient_colors)
                wc_params['colormap'] = custom_cmap
            else:
                # Use preset colormap
                wc_params['colormap'] = self.selected_colormap
            
            # Set background and mode
            if self.rgba_mode.get():
                wc_params['mode'] = 'RGBA'
                wc_params['background_color'] = None
            else:
                wc_params['mode'] = 'RGB'
                wc_params['background_color'] = self.bg_color.get()
            
            # Apply mask based on radio button selection
            mask_to_use = None
            mask_type = self.mask_type.get()
            
            if mask_type == "image_mask" and hasattr(self, 'image_mask_image') and self.image_mask_image is not None:
                mask_to_use = self.image_mask_image
            elif mask_type == "text_mask" and hasattr(self, 'text_mask_image') and self.text_mask_image is not None:
                mask_to_use = self.text_mask_image
            
            if mask_to_use is not None:
                wc_params['mask'] = mask_to_use
                # Disable outlines in RGBA mode due to wordcloud library bug
                # (shape mismatch between RGBA image and RGB outline)
                if self.outline_width.get() > 0 and not self.rgba_mode.get():
                    wc_params['contour_width'] = self.outline_width.get()
                    wc_params['contour_color'] = self.outline_color.get()
                elif self.outline_width.get() > 0 and self.rgba_mode.get():
                    self.print_warning("Outlines disabled in RGBA mode due to library compatibility")
            
            # Use our forbidden words instead of default STOPWORDS
            wc_params['stopwords'] = self.forbidden_words
            
            # Log mask info if using one
            if mask_to_use is not None:
                mask_shape = mask_to_use.shape
                self.print_debug(f"Using mask with shape: {mask_shape}")
                # Count available pixels (black pixels in mask)
                if len(mask_shape) == 3:
                    # Convert to grayscale if RGB
                    gray_mask = np.mean(mask_to_use, axis=2)
                else:
                    gray_mask = mask_to_use
                available_pixels = np.sum(gray_mask < 128)  # Count dark pixels
                total_pixels = mask_shape[0] * mask_shape[1]
                self.print_debug(f"Mask available area: {available_pixels:,} pixels ({available_pixels/total_pixels*100:.1f}% of total)")
            
            self.wordcloud = WordCloud(**wc_params).generate(filtered_text)
            
            # Update UI in main thread
            self.root.after(0, self._update_preview)
            
        except Exception as e:
            error_msg = str(e)
            self.root.after(0, lambda: self.show_toast(f"Error generating word cloud: {error_msg}", "danger"))
        finally:
            self.root.after(0, self._generation_complete)
    
    def _update_preview(self):
        """Update the preview canvas with generated word cloud"""
        self.print_debug("Updating preview with new word cloud")
        
        # Get the word cloud image before clearing (since clear_canvas sets wordcloud to None)
        if not hasattr(self, 'wordcloud') or self.wordcloud is None:
            self.print_fail("No wordcloud object to display")
            return
            
        try:
            wc_image = self.wordcloud.to_image()
        except ValueError as e:
            if "broadcast together with shapes" in str(e):
                self.print_fail("Error: RGBA mode with outlines is not supported due to library limitations")
                self.show_toast("Please disable outlines or switch to RGB mode", "danger")
                return
            else:
                raise
        
        # Ensure preview size is updated
        display_width, display_height = self.calculate_preview_size()
        self.figure.set_size_inches(display_width/100, display_height/100)
        
        # Update canvas widget size
        self.canvas_widget.config(width=display_width, height=display_height)
        
        # Clear canvas but keep the wordcloud object, don't show placeholder
        self.clear_canvas(clear_wordcloud=False, show_placeholder=False)
        
        # Get the current axes (created by clear_canvas)
        ax = self.figure.gca()
        
        if self.rgba_mode.get() and self.show_transparency.get():
            # For RGBA mode with transparency display enabled
            import numpy as np
            
            # Get the actual word cloud dimensions
            wc_height, wc_width = wc_image.size[1], wc_image.size[0]
            
            # Create checkered pattern matching the word cloud size
            checker_size = 20
            checkerboard = np.ones((wc_height, wc_width, 3)) * 0.95  # Light gray base
            
            # Create checker pattern
            for i in range(0, wc_height, checker_size * 2):
                for j in range(0, wc_width, checker_size * 2):
                    checkerboard[i:i+checker_size, j:j+checker_size] = 0.9
                    if i + checker_size < wc_height and j + checker_size < wc_width:
                        checkerboard[i+checker_size:i+2*checker_size, j+checker_size:j+2*checker_size] = 0.9
            
            # Show checkerboard as background
            ax.imshow(checkerboard, extent=[0, wc_width, wc_height, 0], aspect='auto')
            
            # Overlay the word cloud with transparency
            ax.imshow(wc_image, extent=[0, wc_width, wc_height, 0], aspect='auto', interpolation='bilinear')
        elif self.rgba_mode.get():
            # RGBA mode but transparency display disabled - show white background
            ax.set_facecolor('white')
            ax.imshow(wc_image, interpolation='bilinear')
        else:
            # For RGB mode, just show the image
            ax.imshow(wc_image, interpolation='bilinear')
        
        ax.axis('off')
        
        # Reduce padding around the plot
        self.figure.tight_layout(pad=0)
        
        # Add size indicator if preview is scaled down
        actual_width = self.canvas_width.get()
        actual_height = self.canvas_height.get()
        if display_width < actual_width or display_height < actual_height:
            scale_percent = int((display_width / actual_width) * 100)
            reduction = 100 - scale_percent
            #ax.text(0.02, 0.98, f"Preview reduced by {reduction}% to fit\nActual size: {actual_width}×{actual_height}px\nPreview size: {display_width}×{display_height}px", 
            #       transform=ax.transAxes, 
            #       fontsize=9, 
            #       verticalalignment='top',
            #       bbox=dict(boxstyle='round,pad=0.4', facecolor='white', alpha=0.9, edgecolor='gray'))
        
        # Force canvas update
        self.canvas.draw()
        self.canvas.flush_events()
        
        # Ensure the GUI is fully updated
        self.root.update_idletasks()
        
        # Enable save button and show success
        self.save_btn.config(state=NORMAL)
        mode_text = "with transparency" if self.rgba_mode.get() else "with solid background"
        self.show_toast(f"Word cloud generated successfully {mode_text}!", "success")
        
        self.print_debug("Preview updated successfully")
    
    def _generation_complete(self):
        """Called when generation is complete"""
        self.progress.stop()
        self.progress.pack_forget()
        self.generate_btn.config(state=NORMAL)
    
    def save_wordcloud(self):
        """Save generated word cloud"""
        if not hasattr(self, 'wordcloud'):
            self.show_toast("Please generate a word cloud first", "warning")
            return
        
        file_path = filedialog.asksaveasfilename(
            defaultextension=".png",
            filetypes=[
                ("PNG files", "*.png"),
                ("JPEG files", "*.jpg"),
                ("SVG files", "*.svg"),
                ("All files", "*.*")
            ]
        )
        
        if file_path:
            try:
                # Handle different file formats
                if file_path.lower().endswith('.svg'):
                    # For SVG, we need to use a different method
                    import io
                    svg_text = self.wordcloud.to_svg()
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(svg_text)
                else:
                    # For PNG/JPEG
                    if file_path.lower().endswith('.jpg') or file_path.lower().endswith('.jpeg'):
                        if self.rgba_mode.get():
                            self.show_message("JPEG format doesn't support transparency. Image will have white background.", "warning")
                    
                    self.wordcloud.to_file(file_path)
                    
                self.show_message(f"Word cloud saved successfully to: {os.path.basename(file_path)}", "good")
                self.show_toast(f"Word cloud saved successfully!", "success")
            except Exception as e:
                self.show_toast(f"Error saving word cloud: {str(e)}", "danger")

    def get_system_fonts(self):
        """Discover fonts available on the system"""
        fonts = set()
        system = platform.system()
        
        if system == "Windows":
            # Windows font directories
            font_dirs = [
                os.path.join(os.environ.get('WINDIR', 'C:\\Windows'), 'Fonts'),
                os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Microsoft', 'Windows', 'Fonts')
            ]
            
            for font_dir in font_dirs:
                if os.path.exists(font_dir):
                    try:
                        for font_file in os.listdir(font_dir):
                            if font_file.lower().endswith(('.ttf', '.otf')):
                                # Try to extract font name from file
                                font_path = os.path.join(font_dir, font_file)
                                try:
                                    # Try to load and get font name
                                    font = ImageFont.truetype(font_path, 12)
                                    # Use filename without extension as fallback
                                    font_name = os.path.splitext(font_file)[0]
                                    fonts.add(font_name)
                                except:
                                    pass
                    except:
                        pass
            
            # Also try to get fonts from registry (more reliable for font names)
            try:
                import winreg
                reg_path = r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts"
                with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, reg_path) as key:
                    i = 0
                    while True:
                        try:
                            name, value, _ = winreg.EnumValue(key, i)
                            # Extract font name from registry entry
                            font_name = name.split(' (')[0]  # Remove style info
                            fonts.add(font_name)
                            i += 1
                        except WindowsError:
                            break
            except:
                pass
                
        elif system == "Darwin":  # macOS
            font_dirs = [
                "/System/Library/Fonts",
                "/Library/Fonts",
                os.path.expanduser("~/Library/Fonts")
            ]
            
            for font_dir in font_dirs:
                if os.path.exists(font_dir):
                    try:
                        for font_file in os.listdir(font_dir):
                            if font_file.lower().endswith(('.ttf', '.otf', '.ttc')):
                                font_name = os.path.splitext(font_file)[0]
                                fonts.add(font_name)
                    except:
                        pass
                        
        else:  # Linux
            # Try fc-list command
            try:
                result = subprocess.run(['fc-list', ':family'], 
                                      capture_output=True, 
                                      text=True)
                if result.returncode == 0:
                    for line in result.stdout.split('\n'):
                        if line.strip():
                            # Extract font family name
                            font_name = line.split(':')[0].strip()
                            fonts.add(font_name)
            except:
                # Fallback to common font directories
                font_dirs = [
                    "/usr/share/fonts",
                    "/usr/local/share/fonts",
                    os.path.expanduser("~/.fonts")
                ]
                
                for font_dir in font_dirs:
                    if os.path.exists(font_dir):
                        for root, dirs, files in os.walk(font_dir):
                            for font_file in files:
                                if font_file.lower().endswith(('.ttf', '.otf')):
                                    font_name = os.path.splitext(font_file)[0]
                                    fonts.add(font_name)
        
        return sorted(list(fonts))
    
    def validate_fonts(self):
        """Check which fonts are actually available on the system"""
        # Show loading message
        self.root.after(0, lambda: self.show_message("Discovering available fonts...", "info"))
        
        # Get system fonts
        system_fonts = self.get_system_fonts()
        
        # Create a dict of validated fonts
        available = {}
        
        # First, add some guaranteed fallback fonts
        fallback_fonts = {
            "Arial": "Arial",
            "Times New Roman": "Times New Roman",
            "Courier New": "Courier New",
            "Verdana": "Verdana"
        }
        
        # Test common fonts that should work
        common_fonts = {
            "Arial": ["Arial", "arial", "arial.ttf"],
            "Arial Black": ["Arial Black", "ariblk", "ariblk.ttf"],
            "Impact": ["Impact", "impact", "impact.ttf"],
            "Times New Roman": ["Times New Roman", "times", "times.ttf"],
            "Georgia": ["Georgia", "georgia", "georgia.ttf"],
            "Verdana": ["Verdana", "verdana", "verdana.ttf"],
            "Comic Sans MS": ["Comic Sans MS", "comic", "comic.ttf"],
            "Trebuchet MS": ["Trebuchet MS", "trebuc", "trebuc.ttf"],
            "Courier New": ["Courier New", "cour", "cour.ttf"],
            "Calibri": ["Calibri", "calibri", "calibri.ttf"],
            "Cambria": ["Cambria", "cambria", "cambria.ttc"],
            "Tahoma": ["Tahoma", "tahoma", "tahoma.ttf"],
            "Century Gothic": ["Century Gothic", "GOTHIC", "GOTHIC.TTF"],
            "Palatino Linotype": ["Palatino Linotype", "pala", "pala.ttf"],
            "Consolas": ["Consolas", "consola", "consola.ttf"],
            "Segoe UI": ["Segoe UI", "segoeui", "segoeui.ttf"]
        }
        
        # Test each common font
        for display_name, attempts in common_fonts.items():
            for attempt in attempts:
                try:
                    ImageFont.truetype(attempt, 12)
                    available[display_name] = attempt
                    break
                except:
                    continue
        
        # Also check system fonts that match our patterns
        for font in system_fonts:
            # Clean up font name for display
            display_name = font.replace('-', ' ').replace('_', ' ')
            
            # Skip if we already have this font
            if display_name in available:
                continue
                
            # Try to load it
            attempts = [
                font,
                f"{font}.ttf",
                f"{font}.otf",
                font.lower(),
                font.replace(' ', ''),
                font.replace(' ', '-')
            ]
            
            for attempt in attempts:
                try:
                    ImageFont.truetype(attempt, 12)
                    available[display_name] = attempt
                    break
                except:
                    continue
        
        # Update available fonts
        if available:
            self.available_fonts = available
            sorted_fonts = sorted(list(available.keys()))
            self.text_mask_font.set(sorted_fonts[0])
            
            # Update font listbox if it exists (in main thread)
            def update_ui():
                if hasattr(self, 'font_listbox'):
                    self.font_listbox.set_fonts(available)
                self.show_message(f"Found {len(available)} fonts available on your system", "good")
            
            self.root.after(0, update_ui)
    
    def show_toast(self, message, style="info"):
        """Show toast notification"""
        # Print to console based on style
        if style in ["danger", "error"]:
            self.print_fail(message)
        elif style == "warning":
            self.print_warning(message)
        else:
            self.print_info(message)
            
        # Use toast manager for stacked toasts with 15 second timeout
        self.toast_manager.show_toast(message, style, duration=15000)

    def toggle_dark_mode(self):
        """Toggle between dark and light themes"""
        if self.dark_mode.get():
            # Switch to dark themes
            self.themes = self.dark_themes
            # Set to first dark theme if current is light
            if self.current_theme.get() not in self.dark_themes:
                self.current_theme.set(self.dark_themes[0])
        else:
            # Switch to light themes
            self.themes = self.light_themes
            # Set to first light theme if current is dark
            if self.current_theme.get() not in self.light_themes:
                self.current_theme.set(self.light_themes[0])
        
        # Update dropdown values
        if hasattr(self, 'theme_dropdown'):
            self.theme_dropdown['values'] = self.themes
        
        # Apply the theme
        self.change_theme()
        
        # Autosave theme preference
        self.autosave_theme()
    
    def change_theme(self, event=None):
        """Change the application theme"""
        new_theme = self.current_theme.get()
        self.root.style.theme_use(new_theme)
        
        # Update any theme-specific elements
        self.show_toast(f"Theme changed to {new_theme}", "info")
        
        # Update canvas background if needed
        if new_theme in ["darkly", "superhero", "solar", "cyborg", "vapor"]:
            # Dark themes - adjust canvas
            self.figure.patch.set_facecolor('#2b2b2b')
            # Update status bar colors for dark theme
            if hasattr(self, 'status_bar'):
                bg_color = '#1F2937'  # Dark gray
                border_color = '#374151'  # Darker border
                text_color = '#F9FAFB'  # Light text
                label_color = '#9CA3AF'  # Gray labels
                self.update_status_bar_colors(bg_color, border_color, text_color, label_color)
        else:
            # Light themes
            self.figure.patch.set_facecolor('white')
            # Update status bar colors for light theme
            if hasattr(self, 'status_bar'):
                bg_color = '#F3F4F6'  # Light gray
                border_color = '#E5E7EB'  # Light border
                text_color = '#1F2937'  # Dark text
                label_color = '#9CA3AF'  # Gray labels
                self.update_status_bar_colors(bg_color, border_color, text_color, label_color)
        self.canvas.draw()
        
        # Autosave theme preference
        if hasattr(self, 'ui_ready') and self.ui_ready:
            self.autosave_theme()
    
    def autosave_theme(self):
        """Autosave only theme-related settings"""
        try:
            theme_config = {
                'theme': self.current_theme.get(),
                'dark_mode': self.dark_mode.get()
            }
            
            # Save to a separate theme config file
            theme_file = get_resource_path(os.path.join('configs', 'theme.json'))
            
            # Ensure directory exists (especially important for PyInstaller exe)
            theme_dir = os.path.dirname(theme_file)
            if not os.path.exists(theme_dir):
                self.print_info(f"Creating theme config directory: {theme_dir}")
                os.makedirs(theme_dir, exist_ok=True)
            
            with open(theme_file, 'w') as f:
                json.dump(theme_config, f, indent=2)
            
            self.print_debug(f"Theme autosaved to {theme_file}: {theme_config}")
        except Exception as e:
            self.print_warning(f"Could not autosave theme: {e}")
    
    def load_theme_preference(self):
        """Load theme preference from file"""
        try:
            theme_file = get_resource_path(os.path.join('configs', 'theme.json'))
            if os.path.exists(theme_file):
                with open(theme_file, 'r') as f:
                    theme_config = json.load(f)
                
                # Apply dark mode setting
                if 'dark_mode' in theme_config:
                    self.dark_mode.set(theme_config['dark_mode'])
                    if theme_config['dark_mode']:
                        self.themes = self.dark_themes
                    else:
                        self.themes = self.light_themes
                    
                    # Update theme dropdown if it exists
                    if hasattr(self, 'theme_dropdown'):
                        self.theme_dropdown['values'] = self.themes
                
                # Apply theme
                if 'theme' in theme_config and theme_config['theme'] in self.themes:
                    self.current_theme.set(theme_config['theme'])
                    self.root.style.theme_use(theme_config['theme'])
                
                self.print_debug(f"Theme loaded: {theme_config}")
        except Exception as e:
            self.print_debug(f"Could not load theme preference: {e}")
    
    def apply_config(self, config, show_message=True):
        """Apply configuration from dictionary"""
        try:
            self.print_debug("Applying configuration...")
            self.print_debug("Configuration being loaded:")
            for key, value in config.items():
                self.print_debug(f"  {key}: {value}")
            # Apply basic settings
            if 'min_length' in config:
                self.print_debug(f"Loading min_length: {config['min_length']}")
                self.min_word_length.set(config['min_length'])
                if self.min_length_meter:
                    self.min_length_meter.amountusedvar.set(config['min_length'])
                elif self.min_length_scale:
                    self.min_length_scale.set(config['min_length'])
                    self.min_length_label.config(text=str(config['min_length']))
                self.print_debug(f"min_word_length after loading: {self.min_word_length.get()}")
            if 'max_length' in config:
                self.max_word_length.set(config['max_length'])
                if self.max_length_meter:
                    self.max_length_meter.amountusedvar.set(config['max_length'])
                elif self.max_length_scale:
                    self.max_length_scale.set(config['max_length'])
                    self.max_length_label.config(text=str(config['max_length']))
            if 'forbidden_words' in config:
                self.print_debug(f"Loading {len(config['forbidden_words'])} forbidden words from config")
                self.forbidden_text.delete(1.0, tk.END)
                forbidden_text = '\n'.join(config['forbidden_words'])
                self.forbidden_text.insert(1.0, forbidden_text)
                # Update the forbidden words set from the text area
                self.update_forbidden_words(show_toast=False)
                self.print_debug(f"Forbidden words set now has {len(self.forbidden_words)} words")
            
            # Apply color settings
            if 'color_mode' in config:
                self.color_mode.set(config['color_mode'])
                self.on_color_mode_change()
            if 'color_scheme' in config:
                self.color_var.set(config['color_scheme'])
                self.on_color_select()
            if 'single_color' in config:
                self.single_color.set(config['single_color'])
                # Update single color preview
                style = ttk.Style()
                style.configure("SingleColorPreview.TFrame", background=config['single_color'])
                if hasattr(self, 'single_color_preview'):
                    self.single_color_preview.configure(style="SingleColorPreview.TFrame")
            if 'custom_colors' in config:
                self.custom_gradient_colors = config['custom_colors']
                self.update_custom_gradient_preview()
            
            # Apply other settings
            if 'prefer_horizontal' in config:
                pref_val = config['prefer_horizontal'] * 100  # Convert to percentage
                self.horizontal_scale.set(pref_val)
                self.horizontal_gauge.configure(value=pref_val)
            if 'letter_thickness' in config:
                thickness_val = config.get('letter_thickness', 1.0)
                # Clamp to new range (0-5)
                thickness_val = max(0, min(5, thickness_val))
                self.letter_thickness.set(thickness_val)
                if hasattr(self, 'thickness_meter') and self.thickness_meter:
                    self.thickness_meter.amountusedvar.set(thickness_val)
                elif hasattr(self, 'thickness_scale') and self.thickness_scale:
                    self.thickness_scale.set(thickness_val)
            if 'letter_spacing' in config:
                spacing_val = config.get('letter_spacing', 0.0)
                # Clamp to range (0-5)
                spacing_val = max(0, min(5, spacing_val))
                self.letter_spacing.set(spacing_val)
                if hasattr(self, 'spacing_meter') and self.spacing_meter:
                    self.spacing_meter.amountusedvar.set(spacing_val)
            # Canvas settings
            if 'canvas_width' in config:
                self.print_debug(f"Setting canvas width to: {config['canvas_width']}")
                self.canvas_width.set(config['canvas_width'])
                # Update UI elements
                if hasattr(self, 'width_meter') and self.width_meter:
                    self.width_meter.amountusedvar.set(config['canvas_width'])
                elif hasattr(self, 'width_scale') and self.width_scale:
                    self.width_scale.set(config['canvas_width'])
                    if hasattr(self, 'width_label'):
                        self.width_label.config(text=f"{config['canvas_width']} px")
            else:
                self.print_warning("canvas_width not found in config")
                
            if 'canvas_height' in config:
                self.print_debug(f"Setting canvas height to: {config['canvas_height']}")
                self.canvas_height.set(config['canvas_height'])
                # Update UI elements
                if hasattr(self, 'height_meter') and self.height_meter:
                    self.height_meter.amountusedvar.set(config['canvas_height'])
                elif hasattr(self, 'height_scale') and self.height_scale:
                    self.height_scale.set(config['canvas_height'])
                    if hasattr(self, 'height_label'):
                        self.height_label.config(text=f"{config['canvas_height']} px")
            else:
                self.print_warning("canvas_height not found in config")
                
            if 'background_color' in config:
                self.print_debug(f"Setting background color to: {config['background_color']}")
                self.bg_color.set(config['background_color'])
                self.update_bg_preview()
            else:
                self.print_warning("background_color not found in config")
                
            if 'rgba_mode' in config:
                self.print_debug(f"Setting RGBA mode to: {config['rgba_mode']}")
                self.rgba_mode.set(config['rgba_mode'])
                self.on_color_mode_change_canvas()
                self.update_mode(show_toast=False)  # Also update the UI state
            # Legacy support for old config files
            elif 'color_mode_setting' in config:
                self.print_debug(f"Using legacy color_mode_setting: {config['color_mode_setting']}")
                if config['color_mode_setting'] == 'RGBA':
                    self.rgba_mode.set(True)
                else:
                    self.rgba_mode.set(False)
                self.on_color_mode_change_canvas()
                self.update_mode(show_toast=False)  # Also update the UI state
            else:
                self.print_warning("rgba_mode not found in config")
            if 'max_words' in config:
                self.max_words.set(config['max_words'])
                if self.max_words_meter:
                    self.max_words_meter.amountusedvar.set(config['max_words'])
                elif self.max_words_scale:
                    self.max_words_scale.set(config['max_words'])
            if 'scale' in config:
                self.scale.set(config['scale'])
                if self.scale_meter:
                    self.scale_meter.amountusedvar.set(config['scale'])
                elif self.scale_scale:
                    self.scale_scale.set(config['scale'])
            
            # Don't load theme from config - it's handled separately in theme.json
            
            # Apply mask settings
            if 'mask_type' in config:
                mask_types = {'no_mask': 0, 'image_mask': 1, 'text_mask': 2}
                if config['mask_type'] in mask_types:
                    self.mask_notebook.select(mask_types[config['mask_type']])
            
            # Apply outline settings
            if hasattr(self, 'outline_var'):
                if 'outline_enabled' in config:
                    self.outline_var.set(config['outline_enabled'])
                if 'outline_width' in config:
                    self.outline_width_var.set(config['outline_width'])
                if 'outline_color' in config:
                    self.outline_color = config['outline_color']
                    self.update_outline_color_preview()
            
            # Apply text mask settings
            if hasattr(self, 'mask_text_var'):
                if 'text_mask_text' in config:
                    self.mask_text_var.set(config['text_mask_text'])
                if 'text_mask_font' in config:
                    self.selected_font.set(config['text_mask_font'])
                if 'text_mask_size' in config:
                    self.text_size_var.set(config['text_mask_size'])
                if 'text_mask_bold' in config:
                    self.bold_var.set(config['text_mask_bold'])
                if 'text_mask_width' in config:
                    self.text_width_var.set(config['text_mask_width'])
                if 'text_mask_height' in config:
                    self.text_height_var.set(config['text_mask_height'])
                if 'text_mask_lock_aspect' in config and hasattr(self, 'lock_aspect_var'):
                    self.lock_aspect_var.set(config['text_mask_lock_aspect'])
                
                # Update text mask if text is present
                if config.get('text_mask_text'):
                    self.update_text_mask()
            
            # Apply image mask settings
            if 'mask_path' in config and hasattr(self, 'mask_path'):
                self.mask_path.set(config['mask_path'])
                # Try to reload the mask if path exists
                mask_path_value = self.mask_path.get()
                if mask_path_value and os.path.exists(mask_path_value):
                    try:
                        self.mask = np.array(Image.open(mask_path_value))
                        # Update mask preview if it exists
                        if hasattr(self, 'update_mask_preview'):
                            self.update_mask_preview()
                    except:
                        pass
            
            # Apply input settings
            if 'working_directory' in config and hasattr(self, 'working_folder'):
                self.working_folder.set(config['working_directory'])
                if os.path.exists(config['working_directory']):
                    self.populate_file_list(show_toast=False)
                    self.print_debug(f"Populated file list for directory: {config['working_directory']}")
            
            # Load pasted text if present
            if 'pasted_text' in config and hasattr(self, 'text_input'):
                self.text_input.delete('1.0', tk.END)
                self.text_input.insert('1.0', config['pasted_text'])
                self.print_debug(f"Loaded pasted text: {len(config['pasted_text'])} characters")
            
            if show_message:
                self.show_message("Configuration loaded successfully", "good")
            
            self.print_info("Configuration applied successfully")
            return True
            
        except Exception as e:
            self.print_fail(f"Error applying config: {str(e)}")
            import traceback
            self.print_debug(traceback.format_exc())
            if show_message:
                self.show_message(f"Failed to apply config: {str(e)}", "fail")
            return False
    
    def import_config(self):
        """Import configuration from JSON file"""
        file_path = filedialog.askopenfilename(
            title="Import Configuration",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")]
        )
        
        if file_path:
            try:
                with open(file_path, 'r') as f:
                    config = json.load(f)
                self.apply_config(config)
            except Exception as e:
                self.show_message(f"Failed to import config: {str(e)}", "fail")
    
    def auto_load_config(self):
        """Auto-load configuration from local file if it exists"""
        config_loaded = False
        config_file = get_resource_path(os.path.join('configs', 'default.json'))
        if not os.path.exists(config_file):
            # Try configs directory
            config_file = get_resource_path(os.path.join('configs', 'default.json'))
            
        if os.path.exists(config_file):
            try:
                self.print_info(f"Auto-loading configuration from {config_file}")
                with open(config_file, 'r') as f:
                    content = f.read().strip()
                    if content:  # Only parse if file has content
                        config = json.loads(content)
                        self.apply_config(config, show_message=False)
                        self.print_info(f"Auto-loaded configuration from {config_file}")
                        config_loaded = True
                    else:
                        self.print_warning("Config file is empty, skipping auto-load")
            except json.JSONDecodeError as e:
                self.print_fail(f"Invalid JSON in config file: {e}")
                self.print_warning("Consider deleting the config file or fixing the JSON syntax")
            except Exception as e:
                self.print_fail(f"Failed to auto-load config: {e}")
        
        # If no config was loaded, populate forbidden words with defaults
        if not config_loaded:
            self.print_debug("No config loaded, using default forbidden words")
            self.forbidden_text.insert('1.0', self.default_forbidden)
            self.update_forbidden_words(show_toast=False)
    
    def get_current_config(self):
        """Get current configuration as dictionary"""
        config = {}
        
        # Basic settings
        if hasattr(self, 'min_word_length'):
            config['min_length'] = self.min_word_length.get()
        if hasattr(self, 'max_word_length'):
            config['max_length'] = self.max_word_length.get()
        if hasattr(self, 'forbidden_text'):
            # Get forbidden words and filter out empty lines
            forbidden_text = self.forbidden_text.get(1.0, tk.END).strip()
            config['forbidden_words'] = [word.strip() for word in forbidden_text.split('\n') if word.strip()]
        
        # Color settings
        if hasattr(self, 'color_mode'):
            config['color_mode'] = self.color_mode.get()
        if hasattr(self, 'color_var'):
            config['color_scheme'] = self.color_var.get()
        if hasattr(self, 'single_color'):
            config['single_color'] = self.single_color.get()
        if hasattr(self, 'custom_gradient_colors'):
            config['custom_colors'] = self.custom_gradient_colors
        
        # Canvas settings
        if hasattr(self, 'horizontal_scale'):
            config['prefer_horizontal'] = self.horizontal_scale.get() / 100.0  # Convert percentage to 0-1
        if hasattr(self, 'letter_thickness'):
            config['letter_thickness'] = self.letter_thickness.get()
        if hasattr(self, 'letter_spacing'):
            config['letter_spacing'] = self.letter_spacing.get()
        if hasattr(self, 'canvas_width'):
            config['canvas_width'] = self.canvas_width.get()
        if hasattr(self, 'canvas_height'):
            config['canvas_height'] = self.canvas_height.get()
        if hasattr(self, 'bg_color'):
            config['background_color'] = self.bg_color.get()
        if hasattr(self, 'rgba_mode'):
            config['rgba_mode'] = self.rgba_mode.get()
        
        # Other settings
        if hasattr(self, 'max_words'):
            config['max_words'] = self.max_words.get()
        if hasattr(self, 'scale'):
            config['scale'] = self.scale.get()
        # Don't save theme in main config - it's handled separately in theme.json
        
        # Mask settings
        if hasattr(self, 'mask_notebook'):
            config['mask_type'] = self.get_current_mask_type()
        
        # Image mask settings
        if hasattr(self, 'mask_path'):
            config['mask_path'] = self.mask_path.get()
        if hasattr(self, 'outline_var'):
            config['outline_enabled'] = self.outline_var.get()
        if hasattr(self, 'outline_width_var'):
            config['outline_width'] = self.outline_width_var.get()
        if hasattr(self, 'outline_color'):
            config['outline_color'] = self.outline_color
        
        # Text mask settings
        if hasattr(self, 'mask_text_var'):
            config['text_mask_text'] = self.mask_text_var.get()
        if hasattr(self, 'selected_font'):
            config['text_mask_font'] = self.selected_font.get()
        if hasattr(self, 'text_size_var'):
            config['text_mask_size'] = self.text_size_var.get()
        if hasattr(self, 'bold_var'):
            config['text_mask_bold'] = self.bold_var.get()
        if hasattr(self, 'text_width_var'):
            config['text_mask_width'] = self.text_width_var.get()
        if hasattr(self, 'text_height_var'):
            config['text_mask_height'] = self.text_height_var.get()
        if hasattr(self, 'lock_aspect_var'):
            config['text_mask_lock_aspect'] = self.lock_aspect_var.get()
        
        # Input settings
        if hasattr(self, 'working_folder'):
            config['working_directory'] = self.working_folder.get()
        
        # Save pasted text if any
        if hasattr(self, 'text_input'):
            pasted_text = self.text_input.get('1.0', tk.END).strip()
            if pasted_text:
                config['pasted_text'] = pasted_text
        
        # Note: We don't save default_forbidden as it's only for reset functionality
        
        return config
    
    def get_current_mask_type(self):
        """Get the currently selected mask type"""
        try:
            current_tab = self.mask_notebook.index(self.mask_notebook.select())
            return ['no_mask', 'image_mask', 'text_mask'][current_tab]
        except:
            return 'no_mask'
    
    def save_config_to_file(self, file_path):
        """Save configuration to specified file"""
        try:
            self.print_debug(f"Saving configuration to: {file_path}")
            
            # Ensure directory exists (especially important for PyInstaller exe)
            config_dir = os.path.dirname(file_path)
            if config_dir and not os.path.exists(config_dir):
                self.print_info(f"Creating config directory: {config_dir}")
                os.makedirs(config_dir, exist_ok=True)
            
            config = self.get_current_config()
            
            # Debug print the configuration
            self.print_debug("Configuration to save:")
            for key, value in config.items():
                self.print_debug(f"  {key}: {value}")
            
            # Ensure all values are JSON serializable
            serializable_config = {}
            for key, value in config.items():
                if hasattr(value, 'get'):  # If it's a Tkinter variable
                    serializable_config[key] = value.get()
                else:
                    serializable_config[key] = value
            
            with open(file_path, 'w') as f:
                json.dump(serializable_config, f, indent=2)
            
            self.print_info(f"Configuration saved successfully to: {file_path}")
            return True
        except Exception as e:
            self.print_fail(f"Error saving config: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def export_config(self):
        """Export current configuration to JSON file"""
        file_path = filedialog.asksaveasfilename(
            title="Export Configuration",
            defaultextension=".json",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")]
        )
        
        if file_path:
            if self.save_config_to_file(file_path):
                self.show_message("Configuration exported successfully", "good")
            else:
                self.show_message("Failed to export configuration", "fail")
    
    
    def save_config_locally(self):
        """Save configuration to local file with user feedback"""
        config_file = get_resource_path(os.path.join('configs', 'default.json'))
        if self.save_config_to_file(config_file):
            self.show_message(f"Configuration saved to {os.path.basename(config_file)}", "good")
        else:
            self.show_message("Failed to save configuration locally", "fail")
    
    def on_closing(self):
        """Handle application closing"""
        self.root.quit()
    
    def reset_app(self):
        """Reset application to default settings"""
        from tkinter import messagebox
        
        # Confirm reset
        if messagebox.askyesno("Reset Application", "Are you sure you want to reset all settings to defaults?"):
            # Reset filter settings
            self.min_word_length.set(3)
            self.max_word_length.set(30)
            self.forbidden_text.delete(1.0, tk.END)
            self.forbidden_text.insert(1.0, self.default_forbidden)
            self.update_forbidden_words(show_toast=False)
            
            # Reset color settings
            self.color_mode.set("preset")
            self.color_var.set("Viridis")
            self.single_color.set("#0078D4")
            self.custom_gradient_colors = ["#FF0000", "#00FF00", "#0000FF"]
            self.update_custom_gradient_preview()
            self.on_color_mode_change()  # Update UI to reflect preset mode
            
            # Reset canvas settings
            self.horizontal_scale.set(0.9)
            self.width_var.set(800)
            self.height_var.set(600)
            self.bg_color = "white"
            self.color_mode_var.set("RGB")
            
            # Reset other settings
            self.max_words.set(200)
            if self.max_words_meter:
                self.max_words_meter.amountusedvar.set(200)
            elif self.max_words_scale:
                self.max_words_scale.set(200)
            self.scale.set(1)
            if self.scale_meter:
                self.scale_meter.amountusedvar.set(1)
            elif self.scale_scale:
                self.scale_scale.set(1)
            
            # Reset mask settings
            self.mask_notebook.select(0)  # Select "No Mask" tab
            self.mask_path.set("No mask selected")
            self.mask = None
            self.mask_image = None
            if hasattr(self, 'mask_label'):
                self.mask_label.config(text="No mask selected")
            
            # Reset outline settings
            self.outline_width.set(0)
            if hasattr(self, 'outline_width_meter') and self.outline_width_meter:
                # Use the workaround for resetting to 0
                self.outline_width_meter.configure(amountused=0.001)
                self.outline_width_meter.update_idletasks()
                self.root.after(50, lambda: self.outline_width_meter.configure(amountused=0.0))
                self.root.after(100, lambda: self.outline_width_meter.update_idletasks())
            elif hasattr(self, 'outline_width_scale') and self.outline_width_scale:
                self.outline_width_scale.set(0)
            
            if hasattr(self, 'outline_var'):
                self.outline_var.set(False)
                self.outline_width_var.set(3)
                self.outline_color = 'black'
            
            # Reset text mask settings
            if hasattr(self, 'mask_text_var'):
                self.mask_text_var.set("")
                self.selected_font.set("Arial")
                self.text_size_var.set(100)
                self.bold_var.set(False)
                self.text_width_var.set(800)
                self.text_height_var.set(600)
                if hasattr(self, 'lock_aspect_var'):
                    self.lock_aspect_var.set(False)
                # Clear text mask preview
                if hasattr(self, 'text_mask_preview_label'):
                    self.text_mask_preview_label.config(image='', text="Preview will appear here")
            
            # Reset working directory
            self.working_folder.set("No folder selected")
            if hasattr(self, 'file_listbox'):
                self.file_listbox.delete(0, tk.END)
            
            # Clear loaded text
            self.loaded_text = ""
            if hasattr(self, 'loaded_files_label'):
                self.loaded_files_label.config(text="No files loaded")
            
            # Clear text input area
            if hasattr(self, 'text_area'):
                self.text_area.delete(1.0, tk.END)
            
            # Don't clear canvas on reset
            # self.clear_canvas()
            
            # Reset theme to default
            self.current_theme.set("cosmo")
            self.root.style.theme_use("cosmo")
            
            self.show_message("Application reset to defaults", "good")
    
    def load_assets(self):
        """Load SVG assets as placeholder text for now"""
        # For now, we'll just store the icon names
        # In a real implementation, you would convert SVG to PhotoImage
        self.icon_texts = {
            'tab_input': '📁',
            'tab_filter': '🔍',
            'tab_style': '🎨',
            'btn_generate': '🚀',
            'btn_save': '💾',
            'btn_clear': '🗑️',
            'btn_folder': '📂',
            'icon_image_mask': '🖼️',
            'icon_text_mask': '🔤',
            'icon_no_mask': '⬜'
        }
    
    def show_about(self):
        """Show about dialog with version information"""
        about_text = f"""WordCloud Magic v{self.VERSION}

A modern word cloud generator with beautiful UI and powerful features.

© 2025 WordCloud Magic
Built with Python, Tkinter, and ttkbootstrap

For help and documentation, use the Help menu option."""
        
        messagebox.showinfo("About WordCloud Magic", about_text)
    
    def show_help(self):
        """Show help in browser"""
        # Import at the top to catch errors early
        import webbrowser
        
        try:
            import markdown2
        except ImportError:
            self.show_toast("markdown2 module not installed - please install it with: pip install markdown2", "danger")
            return
        
        try:
            # Get base directory - handle PyInstaller bundle
            if hasattr(sys, '_MEIPASS'):
                base_dir = sys._MEIPASS
            else:
                base_dir = os.path.dirname(os.path.abspath(__file__))
            
            # Read the help.md file from templates folder
            help_md_path = os.path.join(base_dir, 'templates', 'help.md')
            
            # Check if help.md exists
            if not os.path.exists(help_md_path):
                self.show_toast("templates/help.md file not found", "danger")
                return
                
            with open(help_md_path, 'r', encoding='utf-8') as f:
                markdown_content = f.read()
            
            # Convert markdown to HTML with extras for better formatting
            html_content = markdown2.markdown(
                markdown_content,
                extras=['tables', 'fenced-code-blocks', 'header-ids', 'toc']
            )
            
            # Read the HTML template
            template_path = os.path.join(base_dir, 'templates', 'help_template.html')
            if not os.path.exists(template_path):
                self.show_toast("templates/help_template.html not found", "danger")
                return
                
            with open(template_path, 'r', encoding='utf-8') as f:
                html_template = f.read()
            
            # Replace the {content} placeholder with the generated HTML
            full_html = html_template.replace('{content}', html_content)
            
            # Create temp directory if it doesn't exist
            temp_dir = os.path.join(base_dir, 'temp')
            if not os.path.exists(temp_dir):
                os.makedirs(temp_dir)
            
            # Create HTML file in local temp directory
            import time
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            temp_filename = f"help_{timestamp}.html"
            temp_path = os.path.join(temp_dir, temp_filename)
            
            with open(temp_path, 'w', encoding='utf-8') as f:
                f.write(full_html)
            
            # Open in default browser
            # Convert path to proper file URL for Windows
            if os.name == 'nt':  # Windows
                file_url = 'file:///' + temp_path.replace('\\', '/')
            else:
                file_url = f'file://{temp_path}'
            
            print(f"Opening help in browser: {file_url}")  # Debug
            webbrowser.open(file_url)
            
            # Schedule deletion of temp file after a delay
            def cleanup():
                try:
                    os.unlink(temp_path)
                except:
                    pass
            
            self.root.after(30000, cleanup)  # Delete after 30 seconds
            
            self.show_toast("Help opened in your browser", "info")
            
        except Exception as e:
            import traceback
            error_msg = f"Error opening help: {str(e)}"
            print(f"Help error: {error_msg}")
            print(traceback.format_exc())
            self.show_toast(error_msg, "danger")


def main():
    # Create the app with a modern theme
    root = ttk.Window(themename="cosmo")
    icon_path = get_resource_path("icon_256.ico")
    ico_path = get_resource_path("icon_256.ico")
    icon = Image.open(icon_path)
    icon = ImageTk.PhotoImage(icon)
    root.iconphoto(True, icon)
    root.iconbitmap(ico_path)
    app = ModernWordCloudApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()